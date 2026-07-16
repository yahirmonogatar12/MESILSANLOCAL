"""Generación segura y privada de Excel/PowerPoint para conversaciones IA."""

from __future__ import annotations

import hashlib
import json
import logging
import os
import re
import uuid
from collections import Counter, defaultdict
from copy import copy
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any

from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from werkzeug.utils import secure_filename

from app.db import get_db_connection

from .ai_reports import compact_report_result, is_identifier_column, run_report
from .ai_store import artifact_root, json_dumps, json_loads, now_local, retention_deadline

logger = logging.getLogger(__name__)

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_FORMULA_PREFIXES = ("=", "+", "-", "@")
_SAFE_TITLE = re.compile(r"[^\w\- .áéíóúÁÉÍÓÚñÑ가-힣]", re.UNICODE)


def _safe_cell(value: Any) -> Any:
    if value is None or isinstance(value, (int, float, bool, datetime)):
        return value
    text = str(value)
    if text.startswith(_FORMULA_PREFIXES):
        return "'" + text
    return text[:32767]


def _safe_title(value: str, fallback: str) -> str:
    cleaned = _SAFE_TITLE.sub("", str(value or "")).strip()
    return cleaned[:120] or fallback


def _preferred_font(language: str) -> str:
    return "Malgun Gothic" if language == "ko" else "Calibri"


def _numeric_columns(rows: list[dict[str, Any]], columns: list[str]) -> list[str]:
    return [
        column
        for column in columns
        if not is_identifier_column(column)
        and any(isinstance(row.get(column), (int, float)) and not isinstance(row.get(column), bool) for row in rows)
    ]


def _categorical_columns(rows: list[dict[str, Any]], columns: list[str]) -> list[str]:
    result = []
    for column in columns:
        values = [row.get(column) for row in rows if row.get(column) not in (None, "")]
        unique = {str(value) for value in values[:1000]}
        if values and 1 < len(unique) <= 30 and not any(isinstance(value, (int, float)) for value in values):
            result.append(column)
    return result


def _preferred_count_category(report: str, categorical: list[str]) -> str | None:
    if not categorical:
        return None
    priorities = {
        "quality_ict": ("linea", "line", "línea", "resultado", "result", "status"),
        "quality_vision": ("linea", "line", "línea", "resultado", "result", "status"),
        "quality_aoi": ("linea", "line", "línea", "resultado", "result", "status"),
        "quality_lqc": ("turno", "linea", "line", "línea", "status"),
    }
    lookup = {str(column).casefold(): column for column in categorical}
    for candidate in priorities.get(report, ()):
        if candidate.casefold() in lookup:
            return lookup[candidate.casefold()]
    return categorical[0]


def _bom_work_groups(result: dict[str, Any], rows: list[dict[str, Any]]) -> list[tuple[str, list[dict[str, Any]]]]:
    """Agrupa una familia BOM por modelo y devuelve el nombre corto del trabajo."""
    if result.get("report") != "bom" or not rows:
        return []
    model_column = next((name for name in ("modelo", "model", "model_code") if any(row.get(name) for row in rows)), None)
    if not model_column:
        return []
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in rows:
        model = str(row.get(model_column) or "").strip()
        if model:
            grouped[model].append(row)
    family = str((result.get("filters") or {}).get("q") or "").strip()
    used: set[str] = set()
    work_groups = []
    for model in sorted(grouped, key=lambda value: [int(part) if part.isdigit() else part for part in re.split(r"(\d+)", value)]):
        suffix = model[len(family):] if family and model.upper().startswith(family.upper()) else ""
        if not suffix:
            ending = re.search(r"(\d{1,4})$", model)
            suffix = ending.group(1)[-2:] if ending else model
        base = re.sub(r"[\\/*?:\[\]]", "_", suffix).strip()[:31] or "BOM"
        sheet_name = base
        counter = 2
        while sheet_name.casefold() in used:
            marker = f"_{counter}"
            sheet_name = f"{base[:31 - len(marker)]}{marker}"
            counter += 1
        used.add(sheet_name.casefold())
        work_groups.append((sheet_name, grouped[model]))
    return work_groups


def _fill_excel_data_sheet(
    ws,
    rows: list[dict[str, Any]],
    columns: list[str],
    *,
    font_name: str,
    navy: str,
    table_name: str,
) -> None:
    ws.freeze_panes = "A2"
    for col_index, column in enumerate(columns, 1):
        cell = ws.cell(1, col_index, column)
        cell.fill = PatternFill("solid", fgColor=navy)
        cell.font = Font(name=font_name, bold=True, color="FFFFFF")
        cell.alignment = Alignment(horizontal="center")
    for row_index, row in enumerate(rows, 2):
        for col_index, column in enumerate(columns, 1):
            ws.cell(row_index, col_index, _safe_cell(row.get(column)))
    for col_index, column in enumerate(columns, 1):
        sample = [len(str(row.get(column) or "")) for row in rows[:300]]
        ws.column_dimensions[ws.cell(1, col_index).column_letter].width = min(
            45, max(12, len(column) + 2, *(sample or [0]))
        )
    if columns and rows:
        end = ws.cell(len(rows) + 1, len(columns)).coordinate
        table = Table(displayName=table_name, ref=f"A1:{end}")
        table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
        ws.add_table(table)
        # La tabla ya incluye su propio autoFilter en tableN.xml. Agregar además
        # un autoFilter de worksheet sobre el mismo rango hace que Excel repare
        # el libro y elimine la tabla completa al abrirlo.


def _add_lqc_excel_summary(
    ws,
    result: dict[str, Any],
    *,
    font_name: str,
    navy: str,
    light: str,
    include_charts: bool,
) -> None:
    summary = result.get("summary") or {}
    shifts = list(summary.get("by_shift") or [])
    hourly = list(summary.get("by_hour") or [])

    start = ws.max_row + 2
    ws.cell(start, 1, "Turno")
    ws.cell(start, 2, "Horario")
    ws.cell(start, 3, "Escaneos")
    ws.cell(start, 4, "Porcentaje")
    for cell in ws[start]:
        if cell.column <= 4:
            cell.fill = PatternFill("solid", fgColor=navy)
            cell.font = Font(name=font_name, bold=True, color="FFFFFF")
    for offset, item in enumerate(shifts, 1):
        ws.cell(start + offset, 1, item.get("label"))
        ws.cell(
            start + offset,
            2,
            f"{item.get('start')}–{item.get('end')}",
        )
        ws.cell(start + offset, 3, int(item.get("count") or 0))
        ws.cell(start + offset, 4, float(item.get("percentage") or 0) / 100)
        ws.cell(start + offset, 4).number_format = "0.00%"

    hourly_start = start + len(shifts) + 3
    ws.cell(hourly_start, 1, "Inicio de hora")
    ws.cell(hourly_start, 2, "Escaneos")
    for column in range(1, 3):
        ws.cell(hourly_start, column).fill = PatternFill("solid", fgColor=light)
        ws.cell(hourly_start, column).font = Font(name=font_name, bold=True)
    for offset, item in enumerate(hourly, 1):
        ws.cell(hourly_start + offset, 1, item.get("slot"))
        ws.cell(hourly_start + offset, 2, int(item.get("count") or 0))

    if include_charts and shifts:
        shift_chart = BarChart()
        shift_chart.title = "Escaneos LQC por turno"
        shift_chart.y_axis.title = "Escaneos"
        shift_chart.x_axis.title = "Turno"
        shift_chart.height = 7.5
        shift_chart.width = 14
        shift_chart.add_data(
            Reference(
                ws,
                min_col=3,
                min_row=start,
                max_row=start + len(shifts),
            ),
            titles_from_data=True,
        )
        shift_chart.set_categories(
            Reference(
                ws,
                min_col=1,
                min_row=start + 1,
                max_row=start + len(shifts),
            )
        )
        ws.add_chart(shift_chart, f"F{start}")

    if include_charts and hourly:
        hourly_chart = LineChart()
        hourly_chart.title = "Actividad LQC por hora"
        hourly_chart.y_axis.title = "Escaneos"
        hourly_chart.x_axis.title = "Hora de inicio"
        hourly_chart.height = 7.5
        hourly_chart.width = 14
        hourly_chart.add_data(
            Reference(
                ws,
                min_col=2,
                min_row=hourly_start,
                max_row=hourly_start + len(hourly),
            ),
            titles_from_data=True,
        )
        hourly_chart.set_categories(
            Reference(
                ws,
                min_col=1,
                min_row=hourly_start + 1,
                max_row=hourly_start + len(hourly),
            )
        )
        ws.add_chart(hourly_chart, f"F{hourly_start}")


def _add_grouped_analysis_excel_summary(
    ws,
    result: dict[str, Any],
    *,
    font_name: str,
    navy: str,
    light: str,
    include_charts: bool,
) -> None:
    summary = result.get("summary") or {}
    ws.append(["Registros fuente", summary.get("total_registros", 0)])
    ws.append(["Cantidad total", summary.get("cantidad_total", 0)])
    ws.append(["Números de parte", summary.get("numeros_parte", 0)])
    ws.append(["Nivel de detalle", summary.get("detail_level", "numero_parte_turno")])

    shift_rows = list(summary.get("by_shift") or [])
    start = ws.max_row + 2
    ws.cell(start, 1, "Turno")
    ws.cell(start, 2, "Registros")
    ws.cell(start, 3, "Cantidad total")
    for column in range(1, 4):
        ws.cell(start, column).fill = PatternFill("solid", fgColor=navy)
        ws.cell(start, column).font = Font(name=font_name, bold=True, color="FFFFFF")
    for offset, item in enumerate(shift_rows, 1):
        ws.cell(start + offset, 1, item.get("label"))
        ws.cell(start + offset, 2, item.get("registros", 0))
        ws.cell(start + offset, 3, item.get("cantidad_total", 0))

    movement_rows = list(summary.get("by_movement") or [])
    movement_start = start + len(shift_rows) + 3
    ws.cell(movement_start, 1, "Tipo de movimiento")
    ws.cell(movement_start, 2, "Registros")
    ws.cell(movement_start, 3, "Cantidad total")
    for column in range(1, 4):
        ws.cell(movement_start, column).fill = PatternFill("solid", fgColor=light)
        ws.cell(movement_start, column).font = Font(name=font_name, bold=True)
    for offset, item in enumerate(movement_rows, 1):
        ws.cell(movement_start + offset, 1, item.get("label"))
        ws.cell(movement_start + offset, 2, item.get("registros", 0))
        ws.cell(movement_start + offset, 3, item.get("cantidad_total", 0))

    top_parts = list(summary.get("top_parts") or [])[:15]
    part_start = movement_start + len(movement_rows) + 3
    ws.cell(part_start, 1, "Número de parte")
    ws.cell(part_start, 2, "Registros")
    ws.cell(part_start, 3, "Cantidad total")
    for column in range(1, 4):
        ws.cell(part_start, column).fill = PatternFill("solid", fgColor=light)
        ws.cell(part_start, column).font = Font(name=font_name, bold=True)
    for offset, item in enumerate(top_parts, 1):
        ws.cell(part_start + offset, 1, item.get("label"))
        ws.cell(part_start + offset, 2, item.get("registros", 0))
        ws.cell(part_start + offset, 3, item.get("cantidad_total", 0))

    if include_charts and shift_rows:
        shift_chart = BarChart()
        shift_chart.title = "Cantidad total por turno"
        shift_chart.y_axis.title = "Cantidad"
        shift_chart.x_axis.title = "Turno"
        shift_chart.height = 7.5
        shift_chart.width = 14
        shift_chart.add_data(
            Reference(ws, min_col=3, min_row=start, max_row=start + len(shift_rows)),
            titles_from_data=True,
        )
        shift_chart.set_categories(
            Reference(ws, min_col=1, min_row=start + 1, max_row=start + len(shift_rows))
        )
        ws.add_chart(shift_chart, f"F{start}")

    if include_charts and top_parts:
        part_chart = BarChart()
        part_chart.title = "Top números de parte por cantidad"
        part_chart.y_axis.title = "Cantidad"
        part_chart.x_axis.title = "Número de parte"
        part_chart.height = 8
        part_chart.width = 16
        part_chart.add_data(
            Reference(ws, min_col=3, min_row=part_start, max_row=part_start + len(top_parts)),
            titles_from_data=True,
        )
        part_chart.set_categories(
            Reference(ws, min_col=1, min_row=part_start + 1, max_row=part_start + len(top_parts))
        )
        ws.add_chart(part_chart, f"F{part_start}")


_PLAN_WEEKDAYS_ES = ("lun", "mar", "mié", "jue", "vie", "sáb", "dom")
_PLAN_MONTHS_ES = (
    "ene", "feb", "mar", "abr", "may", "jun",
    "jul", "ago", "sep", "oct", "nov", "dic",
)


def _plan_date(value: Any) -> date | None:
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value
    text = str(value or "").strip()[:10]
    try:
        return datetime.strptime(text, "%Y-%m-%d").date()
    except ValueError:
        return None


def _plan_date_label(value: date, *, weekday: bool = True) -> str:
    label = value.strftime("%d/%m")
    return f"{label} {_PLAN_WEEKDAYS_ES[value.weekday()]}" if weekday else label


def _plan_template_path() -> Path:
    configured = str(os.getenv("AI_PLAN_PROPOSAL_TEMPLATE") or "").strip()
    if configured:
        return Path(configured).expanduser().resolve()
    return Path(__file__).resolve().parents[3] / "Plan_260716_FINAL.xlsx"


def _plan_block_assignments(rows: list[dict[str, Any]]) -> dict[tuple[str, str], str]:
    """Reconstruye bloques visuales de 9 h cuando el borrador no los persistió."""
    assignments: dict[tuple[str, str], str] = {}
    by_date: dict[str, dict[str, float]] = defaultdict(lambda: defaultdict(float))
    for row in rows:
        by_date[str(row.get("fecha") or "")][str(row.get("linea") or "")] += float(
            row.get("horas") or 0
        )

    def compatible(line: str, existing: list[str]) -> bool:
        line = line.upper()
        other = [item.upper() for item in existing]
        return not (
            (line == "D3" and any(item.startswith("M") for item in other))
            or (line.startswith("M") and "D3" in other)
        )

    for sched_date, line_hours in sorted(by_date.items()):
        blocks: list[dict[str, Any]] = []
        for line, hours in sorted(line_hours.items(), key=lambda item: (-item[1], item[0])):
            candidates = [
                block for block in blocks
                if block["hours"] + hours <= 9.01 and compatible(line, block["lines"])
            ]
            if candidates:
                block = min(candidates, key=lambda item: (item["hours"], item["name"]))
            else:
                block = {
                    "name": f"B{len(blocks) + 1}",
                    "hours": 0.0,
                    "lines": [],
                }
                blocks.append(block)
            block["hours"] += hours
            block["lines"].append(line)
            assignments[(sched_date, line)] = block["name"]
    return assignments


def _build_plan_proposal_excel(
    result: dict[str, Any],
    title: str,
    language: str,
    target: Path,
) -> bool:
    """Genera la propuesta con la misma plantilla visual usada por Planning."""
    template = _plan_template_path()
    if not template.exists():
        logger.warning("Plantilla de propuesta no encontrada: %s", template)
        return False

    wb = load_workbook(template)
    required = {"Resumen", "Horizonte", "Logica", "No planeadas", "Part"}
    plan_sheet_name = next(
        (name for name in wb.sheetnames if name.startswith("Plan ")),
        None,
    )
    if not required.issubset(set(wb.sheetnames)) or not plan_sheet_name:
        logger.warning("La plantilla de propuesta no contiene las seis hojas esperadas")
        return False

    plan_template = wb[plan_sheet_name]

    def find_row(ws, column: int, predicate, fallback: int) -> int:
        for row_no in range(1, ws.max_row + 1):
            if predicate(ws.cell(row_no, column).value):
                return row_no
        return fallback

    subtotal_row = find_row(
        plan_template,
        2,
        lambda value: str(value or "").startswith("Total B"),
        5,
    )
    total_row = find_row(
        plan_template,
        2,
        lambda value: str(value or "").strip().upper() == "TOTAL",
        21,
    )
    footnote_row = find_row(
        plan_template,
        1,
        lambda value: str(value or "").startswith("Inv. antes"),
        23,
    )
    alert_row = find_row(
        wb["No planeadas"],
        2,
        lambda value: str(value or "").startswith("Las partes omitidas"),
        11,
    )
    part_total_row = find_row(
        wb["Part"],
        1,
        lambda value: str(value or "").strip().upper() == "TOTAL PZS",
        42,
    )
    weekend_column = next(
        (
            cell.column
            for cell in wb["Part"][3]
            if "sáb" in str(cell.value or "").lower()
            or "dom" in str(cell.value or "").lower()
        ),
        5,
    )
    green_style = copy(wb["Horizonte"]["C5"]._style)
    # La plantilla validada conserva el estilo amarillo aunque una propuesta
    # concreta no tenga lineas sobre capacidad. En ese caso permanece sin uso.
    yellow_style = (
        copy(wb._cell_styles[24])
        if len(wb._cell_styles) > 24
        else copy(green_style)
    )
    refs = {
        "title": copy(wb["Resumen"]["A1"]._style),
        "section": copy(wb["Resumen"]["B3"]._style),
        "note": copy(wb["Resumen"]["B4"]._style),
        "header": copy(plan_template["A3"]._style),
        "body": copy(plan_template["B4"]._style),
        "body_bold": copy(plan_template["A4"]._style),
        "negative": copy(plan_template["G4"]._style),
        "subtotal": copy(plan_template.cell(subtotal_row, 2)._style),
        "subtotal_green": copy(plan_template.cell(subtotal_row, 6)._style),
        "subtotal_note": copy(plan_template.cell(subtotal_row, 7)._style),
        "total": copy(plan_template.cell(total_row, 2)._style),
        "total_note": copy(plan_template.cell(total_row, 7)._style),
        "footnote": copy(plan_template.cell(footnote_row, 1)._style),
        "green_body": green_style,
        "green_note": copy(wb["Horizonte"]["F5"]._style),
        "yellow_body": yellow_style,
        "step": copy(wb["Logica"]["B3"]._style),
        "alert": copy(wb["No planeadas"].cell(alert_row, 2)._style),
        "weekend_header": copy(wb["Part"].cell(3, weekend_column)._style),
        "weekend_body": copy(wb["Part"].cell(4, weekend_column)._style),
        "part_total_label": copy(wb["Part"].cell(part_total_row, 1)._style),
        "part_total": copy(wb["Part"].cell(part_total_row, 3)._style),
    }
    for sheet in list(wb.worksheets):
        wb.remove(sheet)

    def style(cell, key: str):
        cell._style = copy(refs[key])
        return cell

    def setup(ws, *, landscape: bool = False):
        ws.sheet_view.showGridLines = False
        ws.freeze_panes = "A2"
        ws.page_setup.orientation = "landscape" if landscape else "portrait"
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.page_margins.left = 0.35
        ws.page_margins.right = 0.35
        ws.page_margins.top = 0.5
        ws.page_margins.bottom = 0.5

    rows = list(result.get("rows") or [])
    summary = dict(result.get("summary") or {})
    proposal_id = str(summary.get("proposal_id") or "")
    date_from = _plan_date(summary.get("date_from")) or date.today()
    date_to = _plan_date(summary.get("date_to")) or date_from
    generated_on = _plan_date(summary.get("created_at")) or date.today()
    blocks = _plan_block_assignments(rows)
    block_count = len(set(blocks.values())) or 1

    # Resumen
    ws = wb.create_sheet("Resumen")
    setup(ws)
    widths = {"A": 3, "B": 34, "C": 26, "D": 26, "E": 60, "F": 3}
    for column, width in widths.items():
        ws.column_dimensions[column].width = width
    ws.row_dimensions[1].height = 30
    ws.merge_cells("A1:F1")
    style(ws["A1"], "title").value = (
        f"Plan Proyectado — propuesta de producción ({date_from.strftime('%d/%m/%Y')})"
    )
    style(ws["B3"], "section").value = "Qué es esto"
    ws.merge_cells("B4:E5")
    style(ws["B4"], "note").value = (
        f"Propuesta calculada por el motor {summary.get('engine_version') or 'MES'} "
        f"para {date_from.strftime('%d/%m/%Y')}"
        + (f" a {date_to.strftime('%d/%m/%Y')}" if date_to != date_from else "")
        + ". El Excel conserva la distribución por bloques, líneas, partes, "
        "cantidades, UPH, horas e inventario. Exportarlo no aplica el schedule."
    )
    style(ws["B7"], "section").value = "Resumen general"
    for column, value in enumerate(("Indicador", "Valor", "Indicador", "Valor"), 2):
        style(ws.cell(8, column), "header").value = value
    summary_rows = [
        ("Propuesta ID", proposal_id, "Estado", summary.get("status")),
        ("Fecha del plan", date_from.isoformat(), "Motor", summary.get("engine_version")),
        ("Partes / lotes", summary.get("total_items", len(rows)), "Cantidad total", summary.get("total_qty", 0)),
        ("Horas totales", summary.get("total_hours", 0), "Bloques de 9 h", block_count),
        ("Partes omitidas", summary.get("omitted_count", 0), "Objetivo", summary.get("objective") or "Plan de producción"),
    ]
    for row_no, values in enumerate(summary_rows, 9):
        for column, value in enumerate(values, 2):
            style(ws.cell(row_no, column), "body_bold" if column in (2, 4) else "body").value = _safe_cell(value)
    style(ws["B16"], "section").value = "Carga por línea"
    for column, value in enumerate(("Línea", "Lotes", "Cantidad", "Horas"), 2):
        style(ws.cell(17, column), "header").value = value
    for row_no, item in enumerate(summary.get("by_line") or [], 18):
        values = (item.get("linea"), item.get("lotes"), item.get("cantidad"), item.get("horas"))
        for column, value in enumerate(values, 2):
            style(ws.cell(row_no, column), "green_body" if float(item.get("horas") or 0) <= 9 else "yellow_body").value = value
        ws.cell(row_no, 4).number_format = "#,##0"
        ws.cell(row_no, 5).number_format = "0.00"
    note_row = 19 + len(summary.get("by_line") or [])
    ws.merge_cells(start_row=note_row, start_column=2, end_row=note_row + 1, end_column=5)
    style(ws.cell(note_row, 2), "note").value = (
        f"Se omitieron {int(summary.get('omitted_count') or 0)} partes. "
        "Consulta la hoja 'No planeadas'. La propuesta permanece pendiente hasta "
        "que se confirme expresamente su aplicación al MES."
    )
    ws.print_area = f"A1:F{note_row + 1}"

    # Plan por bloques
    plan_name = f"Plan {date_from.day:02d}-{_PLAN_MONTHS_ES[date_from.month - 1]}"
    ws = wb.create_sheet(plan_name[:31])
    setup(ws, landscape=True)
    for column, width in {
        "A": 7, "B": 7, "C": 16, "D": 10, "E": 8,
        "F": 9, "G": 12, "H": 12, "I": 12,
    }.items():
        ws.column_dimensions[column].width = width
    ws.row_dimensions[1].height = 30
    ws.merge_cells("A1:I1")
    style(ws["A1"], "title").value = (
        f"Plan propuesto — {date_from.strftime('%d/%m/%Y')} "
        f"({block_count} bloques de 9 h; propuesta {proposal_id[:8]})"
    )
    headers = ("Bloque", "Línea", "Parte", "Cantidad", "UPH", "Horas", "Inv. antes", "Inv. después", "Falta el")
    for column, value in enumerate(headers, 1):
        style(ws.cell(3, column), "header").value = value
    grouped: dict[tuple[str, str], list[dict[str, Any]]] = defaultdict(list)
    for item in rows:
        key = (str(item.get("fecha") or ""), blocks.get((str(item.get("fecha") or ""), str(item.get("linea") or "")), "B1"))
        grouped[key].append(item)
    current = 4
    subtotal_rows = []
    for (_sched_date, block), items in sorted(grouped.items()):
        start = current
        lines = sorted({str(item.get("linea") or "") for item in items})
        for item in items:
            values = (
                block, item.get("linea"), item.get("numero_parte"), item.get("cantidad"),
                item.get("uph"), item.get("horas"), item.get("inventario_antes"),
                item.get("inventario_despues"), item.get("fecha_faltante"),
            )
            for column, value in enumerate(values, 1):
                key = "negative" if column == 7 and float(value or 0) < 0 else ("body_bold" if column == 1 else "body")
                style(ws.cell(current, column), key).value = _safe_cell(value)
            ws.cell(current, 4).number_format = "#,##0"
            ws.cell(current, 5).number_format = "#,##0"
            ws.cell(current, 6).number_format = "0.00"
            current += 1
        subtotal = current
        subtotal_rows.append(subtotal)
        style(ws.cell(subtotal, 2), "subtotal").value = f"Total {block} ({'+'.join(lines)})"
        style(ws.cell(subtotal, 4), "subtotal").value = f"=SUM(D{start}:D{subtotal - 1})"
        style(ws.cell(subtotal, 6), "subtotal_green").value = f"=SUM(F{start}:F{subtotal - 1})"
        style(ws.cell(subtotal, 7), "subtotal_note").value = f'=TEXT(F{subtotal}/9,"0%")&" de las 9 h del turno"'
        ws.cell(subtotal, 4).number_format = "#,##0"
        ws.cell(subtotal, 6).number_format = "0.00"
        current += 2
    total_row = current
    style(ws.cell(total_row, 2), "total").value = "TOTAL"
    style(ws.cell(total_row, 4), "total").value = (
        "=" + "+".join(f"D{row}" for row in subtotal_rows) if subtotal_rows else 0
    )
    style(ws.cell(total_row, 6), "total").value = (
        "=" + "+".join(f"F{row}" for row in subtotal_rows) if subtotal_rows else 0
    )
    style(ws.cell(total_row, 7), "total_note").value = f"de {block_count * 9} h ({block_count} bloques x 9 h)"
    ws.cell(total_row, 4).number_format = "#,##0"
    ws.cell(total_row, 6).number_format = "0.00"
    foot_row = total_row + 2
    ws.merge_cells(start_row=foot_row, start_column=1, end_row=foot_row, end_column=9)
    style(ws.cell(foot_row, 1), "footnote").value = (
        "Inv. antes = inventario proyectado el día del faltante. Inv. después = "
        "inventario más cantidad planeada. Cada línea queda completa en un solo bloque; "
        "ningún bloque debe superar 9 h."
    )
    ws.print_area = f"A1:I{foot_row}"

    # Horizonte / capacidad
    ws = wb.create_sheet("Horizonte")
    setup(ws)
    for column, width in {"A": 3, "B": 14, "C": 12, "D": 12, "E": 14, "F": 46, "G": 3}.items():
        ws.column_dimensions[column].width = width
    ws.row_dimensions[1].height = 30
    ws.merge_cells("A1:G1")
    style(ws["A1"], "title").value = "Capacidad y horizonte de la propuesta"
    style(ws["B3"], "section").value = "Carga por línea"
    for column, value in enumerate(("Línea", "Piezas", "Lotes", "Horas", "Diagnóstico"), 2):
        style(ws.cell(4, column), "header").value = value
    for row_no, item in enumerate(summary.get("by_line") or [], 5):
        hours = float(item.get("horas") or 0)
        cell_style = "green_body" if hours <= 9 else "yellow_body"
        values = (
            item.get("linea"), item.get("cantidad"), item.get("lotes"), hours,
            "Dentro de capacidad" if hours <= 9 else "Excede las 9 h",
        )
        for column, value in enumerate(values, 2):
            style(ws.cell(row_no, column), cell_style if column < 6 else ("green_note" if hours <= 9 else "yellow_body")).value = value
        ws.cell(row_no, 3).number_format = "#,##0"
        ws.cell(row_no, 5).number_format = "0.00"
    explanation_row = 7 + len(summary.get("by_line") or [])
    style(ws.cell(explanation_row, 2), "section").value = "Horizonte evaluado"
    ws.merge_cells(start_row=explanation_row + 1, start_column=2, end_row=explanation_row + 3, end_column=6)
    shortage_dates = [_plan_date(item.get("fecha_faltante")) for item in rows]
    shortage_dates = [item for item in shortage_dates if item]
    shortage_text = (
        f"Los lotes evitan faltantes entre {min(shortage_dates).isoformat()} y {max(shortage_dates).isoformat()}. "
        if shortage_dates else ""
    )
    style(ws.cell(explanation_row + 1, 2), "note").value = (
        f"El schedule propuesto cubre {date_from.isoformat()} a {date_to.isoformat()}. "
        + shortage_text
        + "La hoja muestra la capacidad usada por línea; los bloques combinan líneas compatibles hasta 9 h."
    )
    ws.print_area = f"A1:G{explanation_row + 3}"

    # Lógica
    ws = wb.create_sheet("Logica")
    setup(ws)
    for column, width in {"A": 3, "B": 5, "C": 30, "D": 88, "E": 3}.items():
        ws.column_dimensions[column].width = width
    ws.row_dimensions[1].height = 30
    ws.merge_cells("A1:E1")
    style(ws["A1"], "title").value = "Cómo se calcula el plan — paso a paso"
    logic = [
        ("Inventario inicial (I0)", "Suma LGEMM, ISEMM, SVC, DIF, pendiente, rework, SMT e IMD de la referencia semanal."),
        ("Proyección diaria", "I(t) = I(t-1) - P(t) + S(t): consumo de LG menos lo ya surtido o programado."),
        ("Faltante", "Una parte es candidata cuando el inventario proyectado cae debajo de cero dentro de su horizonte."),
        ("Cantidad", "Se cubre el faltante con margen y caja cerrada; el lote busca conservar el remanente objetivo."),
        ("Línea", "Se respeta la línea de ensamble registrada. Una parte no se reasigna a una línea incorrecta."),
        ("Capacidad", "Horas del lote = cantidad / UPH. Las líneas se acomodan en bloques compatibles de hasta 9 h."),
        ("Orden", "Se priorizan faltantes cercanos y se mantienen juntas las familias para reducir cambios de modelo."),
        ("Excepciones", "Faltantes de inventario, CT, UPH, pack, línea o capacidad se reportan; no se inventan datos."),
    ]
    for index, (label, description) in enumerate(logic, 1):
        row_no = index + 2
        ws.row_dimensions[row_no].height = 56
        style(ws.cell(row_no, 2), "step").value = str(index)
        style(ws.cell(row_no, 3), "body_bold").value = label
        style(ws.cell(row_no, 4), "body").value = description
        ws.cell(row_no, 4).alignment = Alignment(vertical="top", wrap_text=True)
    ws.print_area = "A1:E10"

    # No planeadas
    ws = wb.create_sheet("No planeadas")
    setup(ws)
    for column, width in {"A": 3, "B": 46, "C": 12, "D": 62}.items():
        ws.column_dimensions[column].width = width
    ws.row_dimensions[1].height = 30
    ws.merge_cells("A1:D1")
    style(ws["A1"], "title").value = "Partes que NO entraron al plan y por qué"
    for column, value in enumerate(("Motivo", "Partes", "Qué significa / qué hacer"), 2):
        style(ws.cell(3, column), "header").value = value
    exception_counts: Counter[str] = Counter()
    for item in rows:
        try:
            values = json.loads(str(item.get("excepciones_json") or "[]"))
        except (TypeError, ValueError):
            values = []
        exception_counts.update(str(value) for value in values)
    omitted_rows = [
        (
            "Omitidas antes de formar lotes",
            int(summary.get("omitted_count") or 0),
            "Revisar inventario base, CT, UPH, pack, línea activa y capacidad disponible.",
        )
    ]
    explanations = {
        "PARTIAL_CAPACITY": "El lote entró parcialmente por el límite de horas del bloque.",
        "SHORTAGE_REMAINS": "La cantidad parcial no elimina todo el faltante proyectado.",
    }
    omitted_rows.extend(
        (code, count, explanations.get(code, "Revisar la excepción reportada por el motor."))
        for code, count in sorted(exception_counts.items())
    )
    for row_no, values in enumerate(omitted_rows, 4):
        for column, value in enumerate(values, 2):
            style(ws.cell(row_no, column), "body").value = value
    alert_row = 6 + len(omitted_rows)
    ws.merge_cells(start_row=alert_row, start_column=2, end_row=alert_row + 2, end_column=4)
    style(ws.cell(alert_row, 2), "alert").value = (
        "Las partes omitidas no se agregan de manera automática: el motor conserva la "
        "excepción para que Planning corrija el dato o autorice capacidad adicional."
    )
    ws.print_area = f"A1:D{alert_row + 2}"

    # Part: matriz por parte y fecha
    ws = wb.create_sheet("Part")
    setup(ws, landscape=True)
    ws.column_dimensions["A"].width = 16
    ws.column_dimensions["B"].width = 7
    horizon_end = max(date_to, date_from + timedelta(days=9))
    horizon_dates = [date_from + timedelta(days=offset) for offset in range((horizon_end - date_from).days + 1)][:14]
    for index in range(len(horizon_dates)):
        ws.column_dimensions[get_column_letter(index + 3)].width = 10
    last_column = len(horizon_dates) + 2
    last_letter = get_column_letter(last_column)
    ws.row_dimensions[1].height = 30
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=last_column)
    style(ws["A1"], "title").value = "Cómo queda el Part — renglón S (lo que se planea, por fecha)"
    style(ws["A3"], "header").value = "Parte"
    style(ws["B3"], "header").value = "Línea"
    for index, current_date in enumerate(horizon_dates, 3):
        key = "weekend_header" if current_date.weekday() >= 5 else "header"
        style(ws.cell(3, index), key).value = _plan_date_label(current_date)
    parts = sorted({str(item.get("numero_parte") or "") for item in rows})
    line_by_part = {str(item.get("numero_parte") or ""): item.get("linea") for item in rows}
    quantities = defaultdict(int)
    hours = defaultdict(float)
    for item in rows:
        part = str(item.get("numero_parte") or "")
        sched_date = str(item.get("fecha") or "")
        quantities[(part, sched_date)] += int(item.get("cantidad") or 0)
        hours[(str(item.get("linea") or ""), sched_date)] += float(item.get("horas") or 0)
    for row_no, part in enumerate(parts, 4):
        style(ws.cell(row_no, 1), "body").value = _safe_cell(part)
        style(ws.cell(row_no, 2), "body_bold").value = line_by_part.get(part)
        for index, current_date in enumerate(horizon_dates, 3):
            key = "weekend_body" if current_date.weekday() >= 5 else "green_body"
            value = quantities.get((part, current_date.isoformat()))
            style(ws.cell(row_no, index), key).value = value or None
            ws.cell(row_no, index).number_format = "#,##0"
    total_row = 4 + len(parts)
    style(ws.cell(total_row, 1), "part_total_label").value = "TOTAL pzs"
    for index in range(3, last_column + 1):
        style(ws.cell(total_row, index), "part_total").value = f"=SUM({get_column_letter(index)}4:{get_column_letter(index)}{total_row - 1})"
        ws.cell(total_row, index).number_format = "#,##0"
    capacity_title = total_row + 2
    style(ws.cell(capacity_title, 1), "section").value = "Horas por línea (límite 9 h)"
    capacity_header = capacity_title + 1
    style(ws.cell(capacity_header, 1), "header").value = "Línea"
    for index, current_date in enumerate(horizon_dates, 3):
        key = "weekend_header" if current_date.weekday() >= 5 else "header"
        style(ws.cell(capacity_header, index), key).value = _plan_date_label(current_date)
    lines = sorted({str(item.get("linea") or "") for item in rows})
    for row_no, line in enumerate(lines, capacity_header + 1):
        style(ws.cell(row_no, 1), "part_total_label").value = line
        for index, current_date in enumerate(horizon_dates, 3):
            key = "weekend_body" if current_date.weekday() >= 5 else "green_body"
            value = round(hours.get((line, current_date.isoformat()), 0), 2)
            style(ws.cell(row_no, index), key).value = value or None
            ws.cell(row_no, index).number_format = "0.00"
    note_row = capacity_header + len(lines) + 3
    ws.merge_cells(start_row=note_row, start_column=1, end_row=note_row, end_column=last_column)
    style(ws.cell(note_row, 1), "note").value = "Gris = sábado y domingo. Las cantidades aparecen en la fecha propuesta de producción."
    ws.merge_cells(start_row=note_row + 1, start_column=1, end_row=note_row + 1, end_column=last_column)
    style(ws.cell(note_row + 1, 1), "note").value = "Ninguna línea debe pasar de 9 h; las excepciones de capacidad se muestran en 'No planeadas'."
    ws.print_area = f"A1:{last_letter}{note_row + 1}"

    wb.active = 0
    try:
        wb.calculation.fullCalcOnLoad = True
        wb.calculation.forceFullCalc = True
        wb.calculation.calcMode = "auto"
    except AttributeError:
        pass
    target.parent.mkdir(parents=True, exist_ok=True)
    wb.save(target)
    return True


def _build_excel(
    result: dict[str, Any],
    title: str,
    language: str,
    target: Path,
    *,
    include_summary: bool | None = None,
    include_charts: bool | None = None,
) -> None:
    rows = list(result.get("rows") or [])
    columns = list(result.get("columns") or (list(rows[0].keys()) if rows else []))[:60]
    summary = compact_report_result(result, sample_size=10)
    report = str(result.get("report") or "")
    is_bom = report == "bom"
    include_summary = (not is_bom) if include_summary is None else bool(include_summary)
    include_charts = (not is_bom) if include_charts is None else bool(include_charts)
    if report == "plan_proposal":
        if not _build_plan_proposal_excel(result, title, language, target):
            raise RuntimeError(
                "No se pudo cargar la plantilla versionada del plan de produccion; "
                "se cancelo el archivo para no entregar un formato generico."
            )
        return
    if include_charts:
        include_summary = True
    font_name = _preferred_font(language)
    navy = "1F4E79"
    light = "D9EAF7"

    wb = Workbook()
    ws_summary = wb.active
    ws_summary.title = "Resumen"
    ws_summary["A1"] = title
    ws_summary["A1"].font = Font(name=font_name, bold=True, size=18, color="FFFFFF")
    ws_summary["A1"].fill = PatternFill("solid", fgColor=navy)
    ws_summary.merge_cells("A1:D1")
    ws_summary.append([])
    ws_summary.append(["Indicador", "Valor"])
    count_label = {
        "quality_ict": "Pruebas realizadas",
        "quality_lqc": "Escaneos LQC",
    }.get(report, "Registros")
    ws_summary.append([count_label, len(rows)])
    ws_summary.append(["Fuente", result.get("source") or "MES"])
    ws_summary.append(["Generado", now_local().strftime("%Y-%m-%d %H:%M:%S")])
    for cell in ws_summary[3]:
        cell.font = Font(name=font_name, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=navy)
    for key, values in list((summary.get("numeric_summary") or {}).items())[:8]:
        if report != "plan_proposal":
            ws_summary.append([f"{key} - total", values.get("sum")])
    if report == "plan_proposal":
        plan_summary = result.get("summary") or {}
        ws_summary.append(["Propuesta", plan_summary.get("proposal_id")])
        ws_summary.append(["Estado", plan_summary.get("status")])
        ws_summary.append(["Motor", plan_summary.get("engine_version")])
        ws_summary.append(["Fecha inicial", plan_summary.get("date_from")])
        ws_summary.append(["Fecha final", plan_summary.get("date_to")])
        ws_summary.append(["Cantidad total", plan_summary.get("total_qty", 0)])
        ws_summary.append(["Horas totales", plan_summary.get("total_hours", 0)])
        ws_summary.append(["Partes omitidas", plan_summary.get("omitted_count", 0)])
        ws_summary.append([])
        ws_summary.append(["Línea", "Lotes", "Cantidad", "Horas"])
        summary_header = ws_summary.max_row
        for cell in ws_summary[summary_header]:
            cell.font = Font(name=font_name, bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor=navy)
        for item in plan_summary.get("by_line") or []:
            ws_summary.append([
                item.get("linea"), item.get("lotes"),
                item.get("cantidad"), item.get("horas"),
            ])
    if report == "quality_lqc":
        lqc_summary = result.get("summary") or {}
        ws_summary.append(["Jornada operativa", lqc_summary.get("operational_window")])
        ws_summary.append(["Series únicas", lqc_summary.get("unique_serials", 0)])
        ws_summary.append(["Escaneos repetidos", lqc_summary.get("repeated_scans", 0)])
        _add_lqc_excel_summary(
            ws_summary,
            result,
            font_name=font_name,
            navy=navy,
            light=light,
            include_charts=include_charts,
        )
    if report in {"warehouse_analysis", "quality_lqc_analysis"}:
        _add_grouped_analysis_excel_summary(
            ws_summary,
            result,
            font_name=font_name,
            navy=navy,
            light=light,
            include_charts=include_charts,
        )
    ws_summary.column_dimensions["A"].width = 32
    ws_summary.column_dimensions["B"].width = 28

    work_groups = _bom_work_groups(result, rows)
    if work_groups:
        ws_summary.append(["Trabajos", len(work_groups)])
        for index, (sheet_name, work_rows) in enumerate(work_groups, 1):
            _fill_excel_data_sheet(
                wb.create_sheet(sheet_name),
                work_rows,
                columns,
                font_name=font_name,
                navy=navy,
                table_name=f"BOMTrabajo{index}",
            )
    elif report == "warehouse_analysis":
        _fill_excel_data_sheet(
            wb.create_sheet("Analisis por parte"),
            rows,
            columns,
            font_name=font_name,
            navy=navy,
            table_name="AnalisisAlmacen",
        )
        movement_sheets = {
            "entradas": "Entradas",
            "salidas": "Salidas",
            "retornos": "Retornos",
        }
        for index, (movement, sheet_name) in enumerate(movement_sheets.items(), 1):
            movement_rows = [row for row in rows if row.get("tipo_movimiento") == movement]
            if movement_rows:
                _fill_excel_data_sheet(
                    wb.create_sheet(sheet_name),
                    movement_rows,
                    columns,
                    font_name=font_name,
                    navy=navy,
                    table_name=f"AnalisisMovimiento{index}",
                )
    elif report == "quality_lqc_analysis":
        _fill_excel_data_sheet(
            wb.create_sheet("Analisis por parte"),
            rows,
            columns,
            font_name=font_name,
            navy=navy,
            table_name="AnalisisLQC",
        )
    elif report == "plan_proposal":
        _fill_excel_data_sheet(
            wb.create_sheet("Plan propuesto"),
            rows,
            columns,
            font_name=font_name,
            navy=navy,
            table_name="PlanPropuesto",
        )
    else:
        _fill_excel_data_sheet(
            wb.create_sheet("Datos"),
            rows,
            columns,
            font_name=font_name,
            navy=navy,
            table_name="DatosMES",
        )

    ws_criteria = wb.create_sheet("Criterios")
    ws_criteria.append(["Criterio", "Valor"])
    for key, value in (result.get("filters") or {}).items():
        if value not in (None, ""):
            ws_criteria.append([key, _safe_cell(value)])
    ws_criteria.append(["Reporte", result.get("report")])
    ws_criteria.append(["Fuente", result.get("source")])
    ws_criteria.append(["Registros", len(rows)])
    for cell in ws_criteria[1]:
        cell.fill = PatternFill("solid", fgColor=navy)
        cell.font = Font(name=font_name, bold=True, color="FFFFFF")
    ws_criteria.column_dimensions["A"].width = 28
    ws_criteria.column_dimensions["B"].width = 60

    numeric = _numeric_columns(rows, columns)
    categorical = _categorical_columns(rows, columns)
    if (
        include_charts
        and categorical
        and report not in {"quality_lqc", "warehouse_analysis", "quality_lqc_analysis"}
    ):
        category = _preferred_count_category(report, categorical)
        count_tests = report == "quality_ict" or not numeric
        metric = "Pruebas" if count_tests else numeric[0]
        totals: Counter[str] = Counter()
        for row in rows:
            label = str(row.get(category) or "Sin valor")
            if count_tests:
                totals[label] += 1
            else:
                try:
                    totals[label] += float(row.get(metric) or 0)
                except (TypeError, ValueError):
                    continue
        chart_rows = totals.most_common(12)
        start = ws_summary.max_row + 3
        ws_summary.cell(start, 1, category).fill = PatternFill("solid", fgColor=light)
        ws_summary.cell(start, 2, metric).fill = PatternFill("solid", fgColor=light)
        for offset, (label, value) in enumerate(chart_rows, 1):
            ws_summary.cell(start + offset, 1, label)
            ws_summary.cell(start + offset, 2, value)
        if chart_rows:
            chart = BarChart()
            chart.title = f"{metric} por {category}"
            chart.height = 8
            chart.width = 15
            chart.add_data(Reference(ws_summary, min_col=2, min_row=start, max_row=start + len(chart_rows)), titles_from_data=True)
            chart.set_categories(Reference(ws_summary, min_col=1, min_row=start + 1, max_row=start + len(chart_rows)))
            ws_summary.add_chart(chart, f"D{start}")

    if not include_summary:
        wb.remove(ws_summary)

    wb.save(target)


def _add_ppt_text(slide, text: str, left, top, width, height, *, size=20, bold=False, color=(31, 78, 121), font="Calibri"):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)
    return box


def _build_powerpoint(result: dict[str, Any], title: str, language: str, target: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("Falta instalar python-pptx para generar presentaciones") from exc

    rows = list(result.get("rows") or [])
    columns = list(result.get("columns") or (list(rows[0].keys()) if rows else []))
    compact = compact_report_result(result, sample_size=8)
    font = _preferred_font(language)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    logo = Path(__file__).resolve().parents[2] / "static" / "images" / "ilsan-logo.png"

    def slide_title(slide, value: str):
        _add_ppt_text(slide, value, Inches(.6), Inches(.3), Inches(11.8), Inches(.65), size=26, bold=True, font=font)
        if logo.exists():
            slide.shapes.add_picture(str(logo), Inches(11.9), Inches(.18), width=Inches(.9))

    slide = prs.slides.add_slide(blank)
    _add_ppt_text(slide, title, Inches(.8), Inches(2.2), Inches(11.7), Inches(1.2), size=32, bold=True, font=font)
    _add_ppt_text(slide, f"{result.get('title')} · {now_local():%Y-%m-%d}", Inches(.85), Inches(3.5), Inches(11), Inches(.6), size=18, color=(90, 90, 90), font=font)
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(10.8), Inches(5.8), width=Inches(1.5))

    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Alcance / Scope / 범위")
    filters = [f"{key}: {value}" for key, value in (result.get("filters") or {}).items() if value not in (None, "")]
    scope_lines = [f"Reporte: {result.get('title')}", f"Registros: {len(rows)}", f"Fuente: {result.get('source')}"] + filters[:8]
    _add_ppt_text(slide, "\n".join(f"• {line}" for line in scope_lines), Inches(.9), Inches(1.3), Inches(11.5), Inches(5.3), size=20, color=(45, 45, 45), font=font)

    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Indicadores principales / Key KPIs / 주요 지표")
    metrics = compact.get("numeric_summary") or {}
    lines = [f"Registros consultados: {len(rows)}"]
    for key, values in list(metrics.items())[:8]:
        lines.append(f"{key}: total {values.get('sum')} · mín {values.get('min')} · máx {values.get('max')}")
    _add_ppt_text(slide, "\n".join(f"• {line}" for line in lines), Inches(.9), Inches(1.3), Inches(11.4), Inches(5.2), size=20, color=(45, 45, 45), font=font)

    numeric = _numeric_columns(rows, columns)
    categorical = _categorical_columns(rows, columns)
    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Distribución de datos / Data distribution / 데이터 분포")
    chart_added = False
    if numeric and categorical:
        category, metric = categorical[0], numeric[0]
        totals: Counter[str] = Counter()
        for row in rows:
            try:
                totals[str(row.get(category) or "Sin valor")] += float(row.get(metric) or 0)
            except (TypeError, ValueError):
                continue
        points = totals.most_common(10)
        if points:
            data = ChartData()
            data.categories = [label for label, _ in points]
            data.add_series(metric, [value for _, value in points])
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.8), Inches(1.25), Inches(11.8), Inches(5.5), data)
            chart_added = True
    if not chart_added:
        _add_ppt_text(slide, "No hay una combinación categórica y numérica adecuada para una gráfica automática.", Inches(1), Inches(2.4), Inches(11), Inches(1.5), size=22, color=(80, 80, 80), font=font)

    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Hallazgos verificables / Verified findings / 확인된 결과")
    findings = [f"La consulta devolvió {len(rows)} registros."]
    for key, values in list(metrics.items())[:6]:
        findings.append(f"{key}: rango {values.get('min')}–{values.get('max')}; total {values.get('sum')}.")
    if result.get("truncated"):
        findings.append("El resultado alcanzó el límite configurado; se requieren filtros para un análisis completo.")
    _add_ppt_text(slide, "\n".join(f"• {line}" for line in findings), Inches(.9), Inches(1.3), Inches(11.5), Inches(5.2), size=20, color=(45, 45, 45), font=font)

    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Conclusiones / Conclusions / 결론")
    _add_ppt_text(slide, "• Resultados generados exclusivamente con datos consultados del MES.\n• Validar filtros y periodo antes de tomar decisiones operativas.\n• El archivo puede regenerarse para incorporar datos actuales.", Inches(.9), Inches(1.5), Inches(11.3), Inches(4.5), size=22, color=(45, 45, 45), font=font)

    slide = prs.slides.add_slide(blank)
    slide_title(slide, "Fuentes y método / Sources / 출처")
    source_text = (
        f"Fuente MES: {result.get('source')}\n"
        f"Reporte: {result.get('report')}\n"
        f"Filtros: {json.dumps(result.get('filters') or {}, ensure_ascii=False, default=str)}\n"
        f"Registros: {len(rows)}\n"
        f"Generado: {now_local():%Y-%m-%d %H:%M:%S}"
    )
    _add_ppt_text(slide, source_text, Inches(.9), Inches(1.4), Inches(11.4), Inches(4.8), size=19, color=(45, 45, 45), font=font)
    prs.save(target)


def _artifact_record(public_id: str) -> dict[str, Any] | None:
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("SELECT * FROM ai_artifacts WHERE public_id = %s LIMIT 1", (public_id,))
        row = cursor.fetchone()
        return dict(row) if row else None
    finally:
        cursor.close()
        conn.close()


def create_artifact(
    *,
    username: str,
    conversation_id: int,
    message_id: int | None,
    artifact_type: str,
    title: str,
    language: str,
    report_key: str,
    filters: dict[str, Any] | None,
    include_summary: bool | None = None,
    include_charts: bool | None = None,
) -> dict[str, Any]:
    artifact_type = str(artifact_type or "").lower()
    if artifact_type not in {"xlsx", "pptx"}:
        raise ValueError("Formato de artefacto no soportado")
    if artifact_type == "xlsx" and report_key in {
        "quality_lqc", "warehouse_analysis", "quality_lqc_analysis",
    }:
        include_summary = True
        include_charts = True
    if artifact_type == "xlsx" and report_key == "plan_proposal":
        include_summary = True
        include_charts = False
    max_rows = max(1, min(int(os.getenv("AI_ARTIFACT_MAX_ROWS", "10000")), 10000))
    max_bytes = max(1024, int(os.getenv("AI_ARTIFACT_MAX_BYTES", "20971520")))
    result = run_report(username, report_key, filters or {}, limit=max_rows, for_artifact=True)
    if result.get("truncated"):
        raise ValueError("El reporte supera 10,000 filas; agrega filtros más específicos")

    public_id = str(uuid.uuid4())
    safe_title = _safe_title(title, result.get("title") or "Reporte MES")
    timestamp = now_local().strftime("%Y%m%d_%H%M%S")
    filename = secure_filename(f"{safe_title}_{timestamp}.{artifact_type}") or f"reporte_mes_{timestamp}.{artifact_type}"
    target = artifact_root() / f"{public_id}.{artifact_type}"
    try:
        if artifact_type == "xlsx":
            _build_excel(
                result,
                safe_title,
                language,
                target,
                include_summary=include_summary,
                include_charts=include_charts,
            )
            mime = XLSX_MIME
        else:
            _build_powerpoint(result, safe_title, language, target)
            mime = PPTX_MIME
        size = target.stat().st_size
        if size > max_bytes:
            raise ValueError("El archivo generado supera el límite de 20 MB; reduce los filtros")
        digest = hashlib.sha256(target.read_bytes()).hexdigest()
    except Exception:
        target.unlink(missing_ok=True)
        raise

    compact = compact_report_result(result, sample_size=5)
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        cursor.execute(
            """
            INSERT INTO ai_artifacts (
                public_id, conversation_id, message_id, username, artifact_type,
                title, filename, storage_path, mime_type, size_bytes, sha256,
                query_spec_json, filters_json, source_summary_json, row_count,
                language, status, created_at, expires_at
            ) VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,'ready',%s,%s)
            """,
            (
                public_id, conversation_id, message_id, username, artifact_type,
                safe_title, filename, str(target), mime, size, digest,
                json_dumps({
                    "report": report_key,
                    "include_summary": include_summary,
                    "include_charts": include_charts,
                }), json_dumps(filters or {}),
                json_dumps(compact), len(result.get("rows") or []), language,
                now_local(), retention_deadline(),
            ),
        )
        conn.commit()
    except Exception:
        conn.rollback()
        target.unlink(missing_ok=True)
        raise
    finally:
        cursor.close()
        conn.close()

    return public_artifact(_artifact_record(public_id) or {})


def public_artifact(row: dict[str, Any]) -> dict[str, Any]:
    source_summary = json_loads(row.get("source_summary_json"), {}) or {}
    return {
        "id": row.get("public_id"),
        "conversation_id": row.get("conversation_id"),
        "type": row.get("artifact_type"),
        "title": row.get("title"),
        "filename": row.get("filename"),
        "mime_type": row.get("mime_type"),
        "size_bytes": row.get("size_bytes"),
        "sha256": row.get("sha256"),
        "row_count": row.get("row_count"),
        "source": source_summary.get("source"),
        "filters": json_loads(row.get("filters_json"), {}) or {},
        "language": row.get("language"),
        "status": row.get("status"),
        "created_at": row.get("created_at"),
        "expires_at": row.get("expires_at"),
        "download_url": f"/api/ai/artifacts/{row.get('public_id')}/download" if row.get("status") == "ready" else None,
    }


def get_artifact(public_id: str) -> dict[str, Any] | None:
    return _artifact_record(public_id)


def list_artifacts(conversation_id: int) -> list[dict[str, Any]]:
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        cursor.execute(
            "SELECT * FROM ai_artifacts WHERE conversation_id = %s ORDER BY id DESC LIMIT 100",
            (conversation_id,),
        )
        return [public_artifact(dict(row)) for row in (cursor.fetchall() or [])]
    finally:
        cursor.close()
        conn.close()


def regenerate_artifact(
    row: dict[str, Any],
    username: str,
    message_id: int | None = None,
) -> dict[str, Any]:
    query_spec = json_loads(row.get("query_spec_json"), {}) or {}
    filters = json_loads(row.get("filters_json"), {}) or {}
    return create_artifact(
        # Revalidate the report with the permissions of the caller doing the
        # regeneration. This matters when an auditor reviews another user's
        # artifact: audit access must not inherit the owner's source permissions.
        username=username,
        conversation_id=int(row["conversation_id"]),
        message_id=message_id or row.get("message_id"),
        artifact_type=row["artifact_type"],
        title=row["title"],
        language=row.get("language") or "es",
        report_key=query_spec.get("report"),
        filters=filters,
        include_summary=query_spec.get("include_summary"),
        include_charts=query_spec.get("include_charts"),
    )


def artifact_tool_schema(username: str, report_keys: list[str]) -> dict[str, Any] | None:
    if not report_keys:
        return None
    return {
        "type": "function",
        "name": "create_artifact",
        "description": "Crea un Excel profesional o PowerPoint ejecutivo usando un reporte MES autorizado. Para exportar una propuesta pendiente usa plan_proposal con su proposal_id; no uses production_plans porque ese reporte sólo contiene el plan ya aplicado. Para solicitudes de análisis usa warehouse_analysis o quality_lqc_analysis según el módulo; estos Excel incluyen detalle por número de parte/turno, resumen y gráficas. En Excel BOM, include_summary e include_charts deben ser false salvo solicitud expresa.",
        "strict": True,
        "parameters": {
            "type": "object",
            "properties": {
                "artifact_type": {"type": "string", "enum": ["xlsx", "pptx"]},
                "title": {"type": "string", "maxLength": 120},
                "language": {"type": "string", "enum": ["es", "en", "ko"]},
                "report": {"type": "string", "enum": report_keys},
                "include_summary": {"type": ["boolean", "null"]},
                "include_charts": {"type": ["boolean", "null"]},
                "filters": {
                    "type": "object",
                    "properties": {
                        "q": {"type": ["string", "null"]},
                        "part_number": {"type": ["string", "null"]},
                        "serial": {"type": ["string", "null"]},
                        "lot": {"type": ["string", "null"]},
                        "model": {"type": ["string", "null"]},
                        "status": {"type": ["string", "null"]},
                        "date_from": {"type": ["string", "null"]},
                        "date_to": {"type": ["string", "null"]},
                        "language": {"type": ["string", "null"]},
                        "shift": {
                            "type": ["string", "null"],
                            "enum": ["actual", "dia", "tiempo_extra", "noche", None],
                        },
                        "movement_type": {
                            "type": ["string", "null"],
                            "enum": [
                                "entradas", "salidas", "retornos",
                                "entradas_salidas", "todos", None,
                            ],
                        },
                        "proposal_id": {"type": ["string", "null"]},
                    },
                    "required": [
                        "q", "part_number", "serial", "lot", "model", "status",
                        "date_from", "date_to", "language", "shift", "movement_type",
                        "proposal_id",
                    ],
                    "additionalProperties": False,
                },
            },
            "required": [
                "artifact_type", "title", "language", "report", "filters",
                "include_summary", "include_charts",
            ],
            "additionalProperties": False,
        },
    }
