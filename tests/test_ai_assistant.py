import json
from datetime import datetime
from pathlib import Path
from types import SimpleNamespace
from xml.etree import ElementTree
from zipfile import ZipFile

import pytest
from openpyxl import load_workbook
from pptx import Presentation

from app.api.portal import ai_artifacts, ai_assistant, ai_openai, ai_reports, ai_store


def _sample_report():
    return {
        "report": "material_inventory",
        "title": "Inventario de materiales",
        "source": "vista_inventario_consolidado",
        "filters": {"part_number": "KS-100"},
        "columns": ["numero_parte", "cantidad", "estado"],
        "rows": [
            {"numero_parte": "KS-100", "cantidad": 12, "estado": "OK"},
            {"numero_parte": "=HYPERLINK(\"http://evil\")", "cantidad": 5, "estado": "HOLD"},
            {"numero_parte": "KS-101", "cantidad": 8, "estado": "OK"},
        ],
        "row_count": 3,
        "truncated": False,
    }


def test_reporte_bom_usa_vista_canonica_ks_eco(monkeypatch):
    captured = {}

    class Cursor:
        def execute(self, sql, params):
            captured["sql"] = sql
            captured["params"] = params

        def fetchall(self):
            return [{
                "id": 1,
                "modelo": "EBR80757421",
                "codigo_material": "EAH62735001",
                "numero_parte": "EAH62735001",
                "cantidad_total": 1,
                "bom_revision": "R01",
            }]

        def close(self):
            pass

    class Connection:
        def cursor(self):
            return Cursor()

        def close(self):
            pass

    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: Connection())

    result = ai_reports._bom_report({"q": "7421"}, limit=200)

    assert ai_reports.REPORTS["bom"].source == "v_ecos_bom_current"
    assert "FROM v_ecos_bom_current v" in captured["sql"]
    assert "ks_bom_headers + ks_bom_components" in result["source"]
    assert result["rows"][0]["modelo"] == "EBR80757421"
    assert result["filters"]["revision_scope"] == "vigente"
    assert captured["params"][-1] == 201


def test_adjunto_se_informa_al_modelo_y_no_expone_file_ref(tmp_path, monkeypatch):
    monkeypatch.setattr(ai_assistant, "_upload_root", lambda: tmp_path)
    folder = tmp_path / "44"
    folder.mkdir()
    (folder / "abc123.xlsx").write_bytes(b"excel-test")
    (folder / "abc123.name").write_text("Plan LG semana 29.xlsx", "utf-8")

    info = ai_assistant._uploaded_file_info(44, "abc123")
    instructions = ai_openai.build_instructions({
        "language": "es",
        "plan_tools_enabled": True,
        "attachment": info,
    })

    assert info == {
        "filename": "Plan LG semana 29.xlsx",
        "extension": ".xlsx",
        "size_bytes": 10,
        "kind": "excel",
    }
    assert "El archivo ya llegó al servidor" in instructions
    assert "Plan LG semana 29.xlsx" in instructions
    assert "plan_importar_preparar" in instructions
    assert "abc123" not in instructions


def test_adjunto_rechaza_referencia_ajena_o_manipulada(tmp_path, monkeypatch):
    monkeypatch.setattr(ai_assistant, "_upload_root", lambda: tmp_path)
    folder = tmp_path / "44"
    folder.mkdir()
    (folder / "safe.xlsx").write_bytes(b"x")

    assert ai_assistant._uploaded_file_info(44, "../safe") is None
    assert ai_assistant._uploaded_file_info(45, "safe") is None


def test_confirmacion_pendiente_ejecuta_importacion_sin_volver_a_preparar(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    message_ids = iter((71, 72))
    executed = []
    recorded = []
    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(ai_assistant, "get_conversation", lambda _public_id: {
        "id": 44,
        "public_id": "chat-plan",
        "username": "ana",
        "title": "Importar plan",
        "language": "es",
    })
    monkeypatch.setattr(ai_assistant, "check_quota", lambda *_args, **_kwargs: (True, None, {}, {}))
    monkeypatch.setattr(ai_assistant, "get_message_by_client_id", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "add_message", lambda *_args, **_kwargs: next(message_ids))
    monkeypatch.setattr(ai_assistant, "recent_model_messages", lambda *_args, **_kwargs: [])
    monkeypatch.setattr(ai_assistant, "allowed_reports", lambda *_args: [])
    monkeypatch.setattr(ai_assistant, "query_tool_schema", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "_has", lambda *_args: True)
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_args: {})
    monkeypatch.setattr(ai_assistant, "_uploaded_file_info", lambda *_args: None)
    monkeypatch.setattr(ai_assistant.ai_plan_tools, "tool_schemas", lambda *_args: [{"name": "plan"}])
    monkeypatch.setattr(ai_assistant, "get_pending_plan_confirmation", lambda *_args: {
        "prepare_tool": "plan_importar_preparar",
        "execute_tool": "plan_importar_ejecutar",
        "confirm_token": "token-servidor",
    })

    def fake_execute(name, arguments, **_kwargs):
        executed.append((name, arguments))
        return {
            "import_id": 91,
            "plan_partes": 487,
            "plan_fechas": 80,
            "plan_registros": 12000,
            "rango": "2026-06-29 a 2026-09-16",
            "inventario_partes": 404,
            "schedules": 800,
        }

    monkeypatch.setattr(ai_assistant.ai_plan_tools, "execute", fake_execute)
    monkeypatch.setattr(ai_assistant, "record_tool_execution", lambda **kwargs: recorded.append(kwargs))
    monkeypatch.setattr(ai_assistant, "update_message", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "increment_usage", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "refresh_conversation_summary", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "_audit", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(
        ai_assistant,
        "stream_response",
        lambda **_kwargs: (_ for _ in ()).throw(AssertionError("No debe llamar OpenAI")),
    )

    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.post(
        "/api/ai/conversations/chat-plan/messages/stream",
        json={
            "content": "confirmo",
            "client_message_id": "00000000-0000-4000-8000-000000000091",
            "language": "es",
        },
    )

    body = response.get_data(as_text=True)
    assert response.status_code == 200
    assert len(executed) == 1
    assert executed[0] == (
        "plan_importar_ejecutar",
        {"confirm_token": "token-servidor"},
    )
    assert recorded[0]["arguments"]["confirm_token"] == "[redactado]"
    assert "Importación completada correctamente" in body
    assert "487" in body
    assert "vuelve a preparar" not in body


def test_detecta_confirmacion_corta_del_plan():
    for text in ("sí", "si", "confirmo", "yes", "confirm", "확인"):
        assert ai_assistant._PLAN_CONFIRMATION.fullmatch(text)
    assert not ai_assistant._PLAN_CONFIRMATION.fullmatch("si puedes revisa el archivo")


def test_confirmacion_pendiente_conserva_proposal_id_para_exportar(monkeypatch):
    proposal_id = "94508d78-0fc8-41a5-a035-7d2352fbf46f"

    class Cursor:
        def execute(self, _sql, _params):
            pass

        def fetchall(self):
            return [{
                "tool_name": "plan_propuesta_preparar",
                "status": "success",
                "created_at": "2026-07-16 15:33:18",
                "result_summary_json": json.dumps({
                    "confirm_token": "token-servidor",
                    "proposal_id": proposal_id,
                }),
            }]

        def close(self):
            pass

    class Connection:
        def cursor(self):
            return Cursor()

        def close(self):
            pass

    monkeypatch.setattr(ai_store, "get_db_connection", lambda: Connection())

    pending = ai_store.get_pending_plan_confirmation(44, "ana")

    assert pending["prepare_tool"] == "plan_propuesta_preparar"
    assert pending["execute_tool"] == "plan_propuesta_aplicar"
    assert pending["proposal_id"] == proposal_id


def test_confirmacion_pendiente_reconoce_sync_part(monkeypatch):
    class Cursor:
        def execute(self, _sql, _params):
            pass

        def fetchall(self):
            return [{
                "tool_name": "plan_part_sincronizar_preparar",
                "status": "success",
                "created_at": "2026-07-16 17:00:00",
                "result_summary_json": json.dumps({
                    "confirm_token": "token-sync-part",
                    "resumen": {"sheet_name": "Part 16", "schedules": 353},
                }),
            }]

        def close(self):
            pass

    class Connection:
        def cursor(self):
            return Cursor()

        def close(self):
            pass

    monkeypatch.setattr(ai_store, "get_db_connection", lambda: Connection())

    pending = ai_store.get_pending_plan_confirmation(44, "ana")

    assert pending["prepare_tool"] == "plan_part_sincronizar_preparar"
    assert pending["execute_tool"] == "plan_part_sincronizar_ejecutar"
    assert pending["confirm_token"] == "token-sync-part"


def test_texto_confirmacion_sync_part_aclara_alcance():
    text = ai_assistant._plan_completion_text(
        "plan_part_sincronizar_ejecutar",
        {
            "parts": 31,
            "schedules": 353,
            "replaced": 120,
            "scope": "main",
            "excluded_by_scope": 47,
            "skipped_without_active_line": 2,
            "skipped_parts_without_active_line": ["P-SIN-1", "P-SIN-2"],
            "date_from": "2026-07-13",
            "date_to": "2026-09-16",
        },
        "es",
    )

    assert "Schedule del Part sincronizado" in text
    assert "353" in text
    assert "120" in text
    assert "MAIN" in text
    assert "47" in text
    assert "P-SIN-1" in text
    assert "P-SIN-2" in text
    assert "No se modificaron inventario ni plan LG" in text


def test_solicitud_excel_exporta_propuesta_pendiente_sin_aplicarla(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    proposal_id = "94508d78-0fc8-41a5-a035-7d2352fbf46f"
    message_ids = iter((73, 74))
    captured = {}
    artifact = {
        "id": "artifact-plan-1",
        "type": "xlsx",
        "title": "Plan de producción propuesto",
        "filename": "plan_produccion_propuesto.xlsx",
        "status": "ready",
        "row_count": 18,
        "download_url": "/api/ai/artifacts/artifact-plan-1/download",
    }
    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(ai_assistant, "get_conversation", lambda _public_id: {
        "id": 44,
        "public_id": "chat-plan",
        "username": "ana",
        "title": "Plan MAIN",
        "language": "es",
    })
    monkeypatch.setattr(
        ai_assistant,
        "check_quota",
        lambda *_args, **_kwargs: (True, None, {}, {}),
    )
    monkeypatch.setattr(ai_assistant, "get_message_by_client_id", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "add_message", lambda *_args, **_kwargs: next(message_ids))
    monkeypatch.setattr(ai_assistant, "recent_model_messages", lambda *_args, **_kwargs: [])
    monkeypatch.setattr(ai_assistant, "allowed_reports", lambda *_args: [{"key": "plan_proposal"}])
    monkeypatch.setattr(ai_assistant, "query_tool_schema", lambda *_args: None)
    monkeypatch.setattr(
        ai_assistant,
        "artifact_tool_schema",
        lambda *_args: {"type": "function", "name": "create_artifact"},
    )
    monkeypatch.setattr(ai_assistant, "_has", lambda *_args: True)
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_args: {})
    monkeypatch.setattr(ai_assistant, "_uploaded_file_info", lambda *_args: None)
    monkeypatch.setattr(ai_assistant.ai_plan_tools, "tool_schemas", lambda *_args: [{"name": "plan"}])
    monkeypatch.setattr(ai_assistant, "get_pending_plan_confirmation", lambda *_args: {
        "prepare_tool": "plan_propuesta_preparar",
        "execute_tool": "plan_propuesta_aplicar",
        "confirm_token": "token-servidor",
        "proposal_id": proposal_id,
    })

    def fake_create_artifact(**kwargs):
        captured.update(kwargs)
        return artifact

    monkeypatch.setattr(ai_assistant, "create_artifact", fake_create_artifact)
    monkeypatch.setattr(
        ai_assistant.ai_plan_tools,
        "execute",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(
            AssertionError("Exportar no debe aplicar la propuesta")
        ),
    )
    monkeypatch.setattr(
        ai_assistant,
        "stream_response",
        lambda **_kwargs: (_ for _ in ()).throw(AssertionError("No debe llamar OpenAI")),
    )
    monkeypatch.setattr(ai_assistant, "record_tool_execution", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "update_message", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "increment_usage", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "refresh_conversation_summary", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "_audit", lambda *_args, **_kwargs: None)

    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.post(
        "/api/ai/conversations/chat-plan/messages/stream",
        json={
            "content": "pásame el excel",
            "client_message_id": "00000000-0000-4000-8000-000000000092",
            "language": "es",
        },
    )

    body = response.get_data(as_text=True)
    assert response.status_code == 200
    assert captured["report_key"] == "plan_proposal"
    assert captured["filters"] == {"proposal_id": proposal_id}
    assert captured["include_summary"] is True
    assert captured["include_charts"] is False
    assert "event: artifact_ready" in body
    assert "plan_produccion_propuesto.xlsx" in body
    assert "18" in body
    assert "sin aplicarse al MES" in body


def test_nueva_propuesta_adjunta_excel_en_la_misma_respuesta(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    proposal_id = "153df596-3135-4990-b052-1a8ef58528e7"
    message_ids = iter((75, 76))
    plan_calls = []
    captured = {}
    artifact = {
        "id": "artifact-plan-auto-1",
        "type": "xlsx",
        "title": "Plan de producción propuesto - 2026-07-17",
        "filename": "plan_produccion_propuesto_2026-07-17.xlsx",
        "status": "ready",
        "row_count": 18,
        "download_url": "/api/ai/artifacts/artifact-plan-auto-1/download",
    }
    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(ai_assistant, "get_conversation", lambda _public_id: {
        "id": 45,
        "public_id": "chat-plan-auto",
        "username": "ana",
        "title": "Plan MAIN",
        "language": "es",
    })
    monkeypatch.setattr(
        ai_assistant,
        "check_quota",
        lambda *_args, **_kwargs: (True, None, {}, {}),
    )
    monkeypatch.setattr(ai_assistant, "get_message_by_client_id", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "add_message", lambda *_args, **_kwargs: next(message_ids))
    monkeypatch.setattr(ai_assistant, "recent_model_messages", lambda *_args, **_kwargs: [])
    monkeypatch.setattr(ai_assistant, "allowed_reports", lambda *_args: [{"key": "plan_proposal"}])
    monkeypatch.setattr(ai_assistant, "query_tool_schema", lambda *_args: None)
    monkeypatch.setattr(
        ai_assistant,
        "artifact_tool_schema",
        lambda *_args: {"type": "function", "name": "create_artifact"},
    )
    monkeypatch.setattr(ai_assistant, "_has", lambda *_args: True)
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_args: {})
    monkeypatch.setattr(ai_assistant, "_uploaded_file_info", lambda *_args: None)
    monkeypatch.setattr(
        ai_assistant.ai_plan_tools,
        "tool_schemas",
        lambda *_args: [{"type": "function", "name": "plan_propuesta_preparar"}],
    )

    def fake_plan_execute(name, arguments, **_kwargs):
        plan_calls.append((name, arguments))
        assert name == "plan_propuesta_preparar"
        return {
            "proposal_id": proposal_id,
            "confirm_token": "token-servidor",
            "version": 1,
            "engine_version": "mrp-capacity-v3",
            "date_from": "2026-07-17",
            "date_to": "2026-07-17",
            "items": 18,
            "partes": 18,
            "total_qty": 8860,
            "line_summary": [],
            "omitidas_count": 9,
        }

    def fake_create_artifact(**kwargs):
        captured.update(kwargs)
        return artifact

    def fake_stream_response(**kwargs):
        result = kwargs["execute_tool"](
            "plan_propuesta_preparar",
            {
                "fecha_inicio": "2026-07-17",
                "fecha_fin": "2026-07-17",
                "objetivo": "MAIN",
                "proceso_actual": None,
            },
            "call-plan-auto",
        )
        assert result["model_output"]["automatic_artifact"] == artifact
        assert "ya está adjunto" in result["model_output"]["response_policy"]
        yield result["client_event"]
        yield {
            "event": "delta",
            "data": {
                "text": "Propuesta MAIN preparada. También adjunté el Excel; todavía no se aplicó al MES."
            },
        }
        yield {"event": "usage", "data": {"input_tokens": 10, "output_tokens": 5}}

    monkeypatch.setattr(ai_assistant.ai_plan_tools, "execute", fake_plan_execute)
    monkeypatch.setattr(ai_assistant, "create_artifact", fake_create_artifact)
    monkeypatch.setattr(ai_assistant, "stream_response", fake_stream_response)
    monkeypatch.setattr(ai_assistant, "record_tool_execution", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "update_message", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "increment_usage", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "refresh_conversation_summary", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "_audit", lambda *_args, **_kwargs: None)

    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.post(
        "/api/ai/conversations/chat-plan-auto/messages/stream",
        json={
            "content": "Haz el plan de produccion para main para mañana",
            "client_message_id": "00000000-0000-4000-8000-000000000093",
            "language": "es",
        },
    )

    body = response.get_data(as_text=True)
    assert response.status_code == 200
    assert len(plan_calls) == 1
    assert captured["report_key"] == "plan_proposal"
    assert captured["filters"] == {"proposal_id": proposal_id}
    assert captured["include_summary"] is True
    assert captured["include_charts"] is False
    assert "event: artifact_ready" in body
    assert "plan_produccion_propuesto_2026-07-17.xlsx" in body
    assert "todavía no se aplicó al MES" in body


def test_reporte_propuesta_conserva_detalle_y_resumen_por_linea(monkeypatch):
    proposal_id = "94508d78-0fc8-41a5-a035-7d2352fbf46f"

    class Cursor:
        def __init__(self):
            self.sql = ""
            self.executions = []

        def execute(self, sql, params):
            self.sql = " ".join(str(sql).split())
            self.executions.append((self.sql, params))

        def fetchone(self):
            return {
                "id": 91,
                "public_id": proposal_id,
                "version": 1,
                "date_from": "2026-07-17",
                "date_to": "2026-07-17",
                "status": "PENDING_CONFIRMATION",
                "engine_version": "mrp-capacity-v3",
                "total_items": 2,
                "total_qty": 2320,
                "omitted_count": 9,
            }

        def fetchall(self):
            return [
                {
                    "fecha": "2026-07-17",
                    "numero_parte": "EBR80757432",
                    "linea": "M4",
                    "cantidad": 1940,
                    "ct": 12.0,
                    "uph": 300,
                    "horas": 6.47,
                },
                {
                    "fecha": "2026-07-17",
                    "numero_parte": "ACQ91482499",
                    "linea": "D2",
                    "cantidad": 380,
                    "ct": 24.0,
                    "uph": 150,
                    "horas": 2.53,
                },
            ]

        def close(self):
            pass

    class Connection:
        def __init__(self):
            self.cursor_instance = Cursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            pass

    connection = Connection()
    monkeypatch.setattr(ai_reports, "report_allowed", lambda *_args: True)
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)

    result = ai_reports.run_report(
        "ana", "plan_proposal", {"proposal_id": proposal_id}, for_artifact=True
    )

    assert result["columns"] == [
        "fecha", "numero_parte", "linea", "cantidad", "ct", "uph", "horas"
    ]
    assert result["row_count"] == 2
    assert result["summary"]["total_qty"] == 2320
    assert result["summary"]["total_hours"] == 9.0
    assert result["summary"]["by_line"] == [
        {"linea": "D2", "lotes": 1, "cantidad": 380, "horas": 2.53},
        {"linea": "M4", "lotes": 1, "cantidad": 1940, "horas": 6.47},
    ]
    assert connection.cursor_instance.executions[0][1] == (proposal_id, "ana")


def test_excel_propuesta_incluye_plan_y_resumen_por_linea(tmp_path):
    target = tmp_path / "plan_propuesto.xlsx"
    result = {
        "report": "plan_proposal",
        "title": "Propuesta del plan de producción",
        "source": "lg_plan_proposals + lg_plan_proposal_items",
        "filters": {"proposal_id": "94508d78-0fc8-41a5-a035-7d2352fbf46f"},
        "columns": ["fecha", "numero_parte", "linea", "cantidad", "ct", "uph", "horas"],
        "rows": [{
            "fecha": "2026-07-17",
            "numero_parte": "EBR80757432",
            "linea": "M4",
            "cantidad": 1940,
            "ct": 12.0,
            "uph": 300,
            "horas": 6.47,
            "inventario_antes": -626,
            "inventario_despues": 1314,
            "fecha_faltante": "2026-07-20",
            "excepciones_json": "[]",
        }],
        "row_count": 1,
        "truncated": False,
        "summary": {
            "proposal_id": "94508d78-0fc8-41a5-a035-7d2352fbf46f",
            "status": "PENDING_CONFIRMATION",
            "engine_version": "mrp-capacity-v3",
            "date_from": "2026-07-17",
            "date_to": "2026-07-17",
            "total_qty": 1940,
            "total_hours": 6.47,
            "omitted_count": 9,
            "excluded_parts": ["EBR30299365", "EBR30299369"],
            "by_line": [{"linea": "M4", "lotes": 1, "cantidad": 1940, "horas": 6.47}],
        },
    }

    ai_artifacts._build_excel(
        result,
        "Plan de producción propuesto",
        "es",
        target,
        include_summary=True,
        include_charts=False,
    )

    workbook = load_workbook(target, data_only=False)
    assert workbook.sheetnames == [
        "Resumen", "Plan 17-jul", "Horizonte", "Logica", "No planeadas", "Part"
    ]
    plan = workbook["Plan 17-jul"]
    assert [cell.value for cell in plan[3]] == [
        "Bloque", "Línea", "Parte", "Cantidad", "UPH", "Horas",
        "Inv. antes", "Inv. después", "Falta el",
    ]
    assert plan["A4"].value == "B1"
    assert plan["C4"].value == "EBR80757432"
    assert plan["F4"].value == 6.47
    assert plan["G4"].value == -626
    assert plan["D5"].value == "=SUM(D4:D4)"
    assert plan["D7"].value == "=D5"
    assert workbook["Resumen"]["E11"].value == 1940
    assert workbook["Resumen"]["C12"].value == 6.47
    assert workbook["Part"]["A4"].value == "EBR80757432"
    assert workbook["Part"]["C4"].value == 1940
    assert workbook["No planeadas"]["B5"].value == "Excluidas expresamente por Planning"
    assert workbook["No planeadas"]["C5"].value == 2
    assert workbook["No planeadas"]["D5"].value == "EBR30299365, EBR30299369"
    assert all(not sheet.sheet_view.showGridLines for sheet in workbook.worksheets)


def test_excel_profesional_y_formula_injection(tmp_path):
    target = tmp_path / "reporte.xlsx"
    ai_artifacts._build_excel(_sample_report(), "Reporte inventario", "es", target)

    wb = load_workbook(target, data_only=False)
    assert wb.sheetnames == ["Resumen", "Datos", "Criterios"]
    assert wb["Datos"].freeze_panes == "A2"
    assert wb["Datos"].auto_filter.ref is None
    assert next(iter(wb["Datos"].tables.values())).autoFilter.ref
    assert wb["Datos"]["A3"].value.startswith("'")
    assert wb["Resumen"]["B4"].value == 3
    assert wb["Resumen"]._charts


def test_excel_ict_cuenta_pruebas_y_no_suma_identificadores(tmp_path):
    target = tmp_path / "ict.xlsx"
    rows = [
        {"id": 101, "ict": 1, "linea": "M1", "resultado": "OK", "duracion": 10.5},
        {"id": 102, "ict": 1, "linea": "M1", "resultado": "NG", "duracion": 11.0},
        {"id": 103, "ict": 2, "linea": "M2", "resultado": "OK", "duracion": 9.5},
        {"id": 104, "ict": 2, "linea": "M4", "resultado": "OK", "duracion": 12.0},
    ]
    result = {
        "report": "quality_ict",
        "title": "Historial ICT",
        "source": "history_ict",
        "filters": {"date_from": "2026-07-13", "date_to": "2026-07-13"},
        "columns": ["id", "ict", "linea", "resultado", "duracion"],
        "rows": rows,
        "row_count": len(rows),
        "truncated": False,
    }

    ai_artifacts._build_excel(result, "Historial ICT - pruebas de hoy", "es", target)
    wb = load_workbook(target, data_only=False)
    summary = wb["Resumen"]
    indicators = {summary.cell(row, 1).value: summary.cell(row, 2).value for row in range(1, summary.max_row + 1)}

    assert indicators["Pruebas realizadas"] == 4
    assert "id - total" not in indicators
    assert "ict - total" not in indicators
    assert indicators["duracion - total"] == 43
    chart_header_row = next(
        row for row in range(1, summary.max_row + 1)
        if summary.cell(row, 1).value == "linea" and summary.cell(row, 2).value == "Pruebas"
    )
    counts = {
        summary.cell(row, 1).value: summary.cell(row, 2).value
        for row in range(chart_header_row + 1, summary.max_row + 1)
        if summary.cell(row, 1).value
    }
    assert counts == {"M1": 2, "M2": 1, "M4": 1}
    assert summary._charts


def test_excel_lqc_incluye_horarios_y_graficas_por_turno_y_hora(tmp_path):
    rows = [
        {"serial": "S-1", "last_scan": "2026-07-13 08:10:00", "status": "OK"},
        {"serial": "S-2", "last_scan": "2026-07-13 18:20:00", "status": "OK"},
        {"serial": "S-3", "last_scan": "2026-07-13 23:15:00", "status": "OK"},
        {"serial": "S-3", "last_scan": "2026-07-14 02:10:00", "status": "OK"},
    ]
    lqc_summary = ai_reports._quality_lqc_row_context(rows)
    result = {
        "report": "quality_lqc",
        "title": "Historial de liberación LQC",
        "source": "box_scans",
        "filters": {
            "date_scope": "today",
            "operational_date": "2026-07-13",
            "date_from": "2026-07-13 07:30:00",
            "date_to": "2026-07-14 07:29:59",
        },
        "columns": ["serial", "last_scan", "status", "fecha_operativa", "turno"],
        "rows": rows,
        "row_count": len(rows),
        "truncated": False,
        "summary": lqc_summary,
    }
    target = tmp_path / "historial_lqc.xlsx"

    ai_artifacts._build_excel(
        result,
        "Historial LQC de hoy",
        "es",
        target,
        include_summary=True,
        include_charts=True,
    )

    wb = load_workbook(target, data_only=False)
    summary = wb["Resumen"]
    values = {
        str(summary.cell(row, column).value or "")
        for row in range(1, summary.max_row + 1)
        for column in range(1, min(summary.max_column, 4) + 1)
    }
    assert len(summary._charts) == 2
    assert "07:30–17:30" in values
    assert "17:30–22:00" in values
    assert "22:00–07:30 del día siguiente" in values
    assert lqc_summary["total_scans"] == 4
    assert lqc_summary["unique_serials"] == 3
    assert lqc_summary["repeated_scans"] == 1
    assert [item["count"] for item in lqc_summary["by_shift"]] == [1, 1, 2]
    assert [wb["Datos"].cell(row, 5).value for row in range(2, 6)] == [
        "DIA",
        "TIEMPO EXTRA",
        "NOCHE",
        "NOCHE",
    ]


def test_resumen_compacto_no_suma_identificadores_numericos():
    result = {
        "report": "quality_ict",
        "columns": ["id", "ict", "barcode", "duracion"],
        "rows": [
            {"id": 100, "ict": 1, "barcode": 12345, "duracion": 8},
            {"id": 101, "ict": 2, "barcode": 67890, "duracion": 9},
        ],
        "row_count": 2,
    }

    numeric = ai_reports.compact_report_result(result)["numeric_summary"]
    assert set(numeric) == {"duracion"}
    assert numeric["duracion"]["sum"] == 17


def test_excel_bom_familia_crea_hoja_por_trabajo(tmp_path):
    rows = []
    for work in ("11", "12", "17", "22"):
        rows.extend([
            {
                "id": int(f"10{work}"),
                "modelo": f"EBR807574{work}",
                "codigo_material": f"MAT-{work}-A",
                "numero_parte": f"PART-{work}-A",
                "cantidad_total": 1,
            },
            {
                "id": int(f"20{work}"),
                "modelo": f"EBR807574{work}",
                "codigo_material": f"MAT-{work}-B",
                "numero_parte": f"PART-{work}-B",
                "cantidad_total": 2,
            },
        ])
    result = {
        "report": "bom",
        "title": "BOM",
        "source": "bom",
        "filters": {"q": "EBR807574"},
        "columns": ["id", "modelo", "codigo_material", "numero_parte", "cantidad_total"],
        "rows": rows,
        "row_count": len(rows),
        "truncated": False,
    }
    target = tmp_path / "bom_familia.xlsx"

    ai_artifacts._build_excel(result, "BOM familia EBR807574", "es", target)

    wb = load_workbook(target, data_only=False)
    assert wb.sheetnames == ["11", "12", "17", "22", "Criterios"]
    assert "Datos" not in wb.sheetnames
    assert "Resumen" not in wb.sheetnames
    for work in ("11", "12", "17", "22"):
        ws = wb[work]
        assert ws.freeze_panes == "A2"
        assert ws.auto_filter.ref is None
        assert ws.max_row == 3
        assert {ws.cell(row, 2).value for row in (2, 3)} == {f"EBR807574{work}"}
        assert len(ws.tables) == 1
        assert next(iter(ws.tables.values())).autoFilter.ref
    assert not any(ws._charts for ws in wb.worksheets)

    namespace = {"x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with ZipFile(target) as package:
        worksheet_parts = [name for name in package.namelist() if name.startswith("xl/worksheets/sheet") and name.endswith(".xml")]
        table_parts = [name for name in package.namelist() if name.startswith("xl/tables/table") and name.endswith(".xml")]
        assert all(
            ElementTree.fromstring(package.read(name)).find("x:autoFilter", namespace) is None
            for name in worksheet_parts
        )
        assert table_parts
        assert all(
            ElementTree.fromstring(package.read(name)).find("x:autoFilter", namespace) is not None
            for name in table_parts
        )


def test_excel_bom_incluye_resumen_y_grafica_solo_si_se_solicitan(tmp_path):
    result = {
        "report": "bom",
        "title": "BOM",
        "source": "bom",
        "filters": {"q": "EBR807574"},
        "columns": ["modelo", "cantidad_total"],
        "rows": [
            {"modelo": "EBR80757411", "cantidad_total": 2},
            {"modelo": "EBR80757412", "cantidad_total": 3},
        ],
        "row_count": 2,
        "truncated": False,
    }
    target = tmp_path / "bom_con_resumen.xlsx"

    ai_artifacts._build_excel(
        result,
        "BOM con resumen",
        "es",
        target,
        include_summary=True,
        include_charts=True,
    )

    wb = load_workbook(target, data_only=False)
    assert wb.sheetnames == ["Resumen", "11", "12", "Criterios"]
    assert wb["Resumen"]._charts


def test_powerpoint_profesional_sin_relaciones_externas(tmp_path):
    target = tmp_path / "reporte.pptx"
    ai_artifacts._build_powerpoint(_sample_report(), "재고 보고서", "ko", target)

    prs = Presentation(target)
    assert 6 <= len(prs.slides) <= 10
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "재고 보고서" in all_text
    assert "vista_inventario_consolidado" in all_text
    for slide in prs.slides:
        for rel in slide.part.rels.values():
            assert not rel.is_external


def test_safe_cell_neutraliza_prefijos_de_formula():
    for value in ("=1+1", "+cmd", "-2+3", "@SUM(A1:A2)"):
        assert ai_artifacts._safe_cell(value).startswith("'")
    assert ai_artifacts._safe_cell(42) == 42


def test_costo_estimado_usa_tarifas_configurables(monkeypatch):
    monkeypatch.setenv("AI_INPUT_COST_PER_MILLION", "2.50")
    monkeypatch.setenv("AI_OUTPUT_COST_PER_MILLION", "10")
    assert str(ai_store.estimated_cost_usd(1_000_000, 500_000)) == "7.500000"


def test_artifact_publico_incluye_fuente_y_filtros():
    artifact = ai_artifacts.public_artifact({
        "public_id": "artifact-1",
        "status": "ready",
        "source_summary_json": json.dumps({"source": "eco_ecn"}),
        "filters_json": json.dumps({"part_number": "KS-100"}),
    })
    assert artifact["source"] == "eco_ecn"
    assert artifact["filters"] == {"part_number": "KS-100"}


@pytest.mark.parametrize(
    ("message", "expected"),
    [
        ("Puedes darme el BOM del 7421", {"q": "7421"}),
        ("Genera el BOM completo del 7421", {"q": "7421"}),
        ("Pon en Excel todos los BOM de la familia EBR807574", {"q": "EBR807574"}),
        ("Show me the BOM for EBR80757421", {"q": "EBR80757421"}),
        ("EBR80757421 BOM 주세요", {"q": "EBR80757421"}),
        ("¿Qué es un BOM?", None),
        ("Explícame el BOM", None),
        ("Dame el BOM", None),
    ],
)
def test_detecta_bom_concreto_para_excel_automatico(message, expected):
    assert ai_assistant._automatic_bom_filters(message) == expected


@pytest.mark.parametrize(
    ("message", "expected"),
    [
        ("Dame el BOM del 7421", {"include_summary": False, "include_charts": False}),
        ("Dame el BOM del 7421 con resumen", {"include_summary": True, "include_charts": False}),
        ("BOM 7421 con resumen y gráficas", {"include_summary": True, "include_charts": True}),
        ("BOM 7421 sin resumen ni gráficas", {"include_summary": False, "include_charts": False}),
        ("BOM 7421 with a chart", {"include_summary": False, "include_charts": True}),
    ],
)
def test_opciones_analiticas_bom_solo_por_solicitud_explicita(message, expected):
    assert ai_assistant._bom_excel_options(message) == expected


@pytest.mark.parametrize(
    ("message", "expected"),
    [
        ("Dame UPH y CT de todos los numeros de parte", True),
        ("Necesito el listado completo de inventario", True),
        ("Give me all work orders", True),
        ("전체 모델의 UPH를 보여줘", True),
        ("Dame el UPH del 7421", False),
    ],
)
def test_detecta_solicitud_de_muchos_datos(message, expected):
    assert ai_assistant._large_data_request(message) is expected


def test_solicitud_bom_genera_excel_automatico(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    artifact = {
        "id": "artifact-bom-1",
        "type": "xlsx",
        "title": "BOM 7421",
        "filename": "bom_7421.xlsx",
        "status": "ready",
        "row_count": 101,
        "download_url": "/api/ai/artifacts/artifact-bom-1/download",
    }
    captured = {}
    message_ids = iter((41, 42))

    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(
        ai_assistant,
        "get_conversation",
        lambda _public_id: {
            "id": 9,
            "public_id": "chat-1",
            "username": "ana",
            "title": "BOM",
            "language": "es",
        },
    )
    monkeypatch.setattr(ai_assistant, "check_quota", lambda *_args, **_kwargs: (True, None, {}, {}))
    monkeypatch.setattr(ai_assistant, "get_message_by_client_id", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "add_message", lambda *_args, **_kwargs: next(message_ids))
    monkeypatch.setattr(ai_assistant, "recent_model_messages", lambda *_args, **_kwargs: [])
    monkeypatch.setattr(ai_assistant, "allowed_reports", lambda *_args: [{"key": "bom"}])
    monkeypatch.setattr(ai_assistant, "query_tool_schema", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "_has", lambda *_args: True)
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_args: {})
    monkeypatch.setattr(ai_assistant.ai_plan_tools, "tool_schemas", lambda *_args: [])

    def fake_create_artifact(**kwargs):
        captured.update(kwargs)
        return artifact

    def fake_stream_response(**kwargs):
        assert kwargs["context"]["automatic_artifact"] == artifact
        assert kwargs["tools"] == []
        yield {"event": "delta", "data": {"text": "Adjunté el Excel completo del BOM."}}
        yield {"event": "usage", "data": {"input_tokens": 10, "output_tokens": 5}}

    monkeypatch.setattr(ai_assistant, "create_artifact", fake_create_artifact)
    monkeypatch.setattr(ai_assistant, "stream_response", fake_stream_response)
    monkeypatch.setattr(ai_assistant, "increment_usage", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "record_tool_execution", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "update_message", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "refresh_conversation_summary", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "_audit", lambda *_args, **_kwargs: None)

    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.post(
        "/api/ai/conversations/chat-1/messages/stream",
        json={
            "content": "Puedes darme el BOM del 7421",
            "client_message_id": "00000000-0000-4000-8000-000000000742",
            "language": "es",
        },
    )

    body = response.get_data(as_text=True)
    assert response.status_code == 200
    assert captured["report_key"] == "bom"
    assert captured["artifact_type"] == "xlsx"
    assert captured["filters"] == {"q": "7421"}
    assert captured["include_summary"] is False
    assert captured["include_charts"] is False
    assert "event: artifact_ready" in body
    assert "bom_7421.xlsx" in body


def test_consulta_masiva_genera_excel_y_no_entrega_tabla_larga(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    artifact = {
        "id": "artifact-raw-1",
        "type": "xlsx",
        "title": "Datos maestros RAW completo",
        "filename": "raw_ct_uph_completo.xlsx",
        "status": "ready",
        "row_count": 387,
        "download_url": "/api/ai/artifacts/artifact-raw-1/download",
    }
    captured = {}
    message_ids = iter((51, 52))

    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(
        ai_assistant,
        "get_conversation",
        lambda _public_id: {
            "id": 10,
            "public_id": "chat-raw",
            "username": "ana",
            "title": "UPH y CT",
            "language": "es",
        },
    )
    monkeypatch.setattr(ai_assistant, "check_quota", lambda *_args, **_kwargs: (True, None, {}, {}))
    monkeypatch.setattr(ai_assistant, "get_message_by_client_id", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "add_message", lambda *_args, **_kwargs: next(message_ids))
    monkeypatch.setattr(ai_assistant, "recent_model_messages", lambda *_args, **_kwargs: [])
    monkeypatch.setattr(
        ai_assistant,
        "allowed_reports",
        lambda *_args: [{"key": "raw_model_standards"}],
    )
    monkeypatch.setattr(
        ai_assistant,
        "query_tool_schema",
        lambda *_args: {"type": "function", "name": "query_mes_report"},
    )
    monkeypatch.setattr(ai_assistant, "artifact_tool_schema", lambda *_args: None)
    monkeypatch.setattr(ai_assistant, "_has", lambda *_args: True)
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_args: {})
    monkeypatch.setattr(ai_assistant.ai_plan_tools, "tool_schemas", lambda *_args: [])
    monkeypatch.setattr(
        ai_assistant,
        "run_report",
        lambda *_args, **_kwargs: {
            "report": "raw_model_standards",
            "title": "Datos maestros RAW: CT y UPH",
            "source": "raw",
            "columns": ["part_no", "model", "c_t", "uph"],
            "rows": [
                {"part_no": f"P-{index}", "model": f"M-{index}", "c_t": 13, "uph": 276}
                for index in range(20)
            ],
            "row_count": 200,
            "truncated": True,
            "filters": {},
        },
    )

    def fake_create_artifact(**kwargs):
        captured.update(kwargs)
        return artifact

    def fake_stream_response(**kwargs):
        result = kwargs["execute_tool"](
            "query_mes_report",
            {"report": "raw_model_standards", "filters": {}},
            "call-raw",
        )
        assert len(result["model_output"]["sample"]) == 8
        assert result["model_output"]["automatic_artifact"] == artifact
        assert "no muestres tablas" in result["model_output"]["response_policy"]
        yield result["client_event"]
        yield {"event": "delta", "data": {"text": "Adjunté el Excel completo con UPH y CT."}}
        yield {"event": "usage", "data": {"input_tokens": 10, "output_tokens": 5}}

    monkeypatch.setattr(ai_assistant, "create_artifact", fake_create_artifact)
    monkeypatch.setattr(ai_assistant, "stream_response", fake_stream_response)
    monkeypatch.setattr(ai_assistant, "increment_usage", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "record_tool_execution", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "update_message", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "refresh_conversation_summary", lambda *_args, **_kwargs: None)
    monkeypatch.setattr(ai_assistant, "_audit", lambda *_args, **_kwargs: None)

    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.post(
        "/api/ai/conversations/chat-raw/messages/stream",
        json={
            "content": "Dame UPH y CT de todos los numeros de parte",
            "client_message_id": "00000000-0000-4000-8000-000000000843",
            "language": "es",
        },
    )

    body = response.get_data(as_text=True)
    assert response.status_code == 200
    assert captured["report_key"] == "raw_model_standards"
    assert captured["filters"] == {}
    assert captured["include_summary"] is False
    assert captured["include_charts"] is False
    assert "event: artifact_ready" in body
    assert "raw_ct_uph_completo.xlsx" in body


def test_reporte_desconocido_y_sin_permiso(monkeypatch):
    with pytest.raises(ValueError, match="no registrado"):
        ai_reports.run_report("ana", "sql_libre", {})

    monkeypatch.setattr(ai_reports, "report_allowed", lambda *_: False)
    with pytest.raises(PermissionError):
        ai_reports.run_report("ana", "bom", {})


def test_lqc_sin_fecha_usa_jornada_operativa_de_hoy(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.sql = ""
            self.executions = []

        def execute(self, sql, params=None):
            self.sql = sql
            self.executions.append((sql, params))

        def fetchall(self):
            if self.sql.startswith("SHOW COLUMNS"):
                return [
                    {"Field": "id", "Type": "int"},
                    {"Field": "last_scan", "Type": "datetime"},
                    {"Field": "barcode", "Type": "varchar(80)"},
                    {"Field": "serial", "Type": "varchar(80)"},
                ]
            return [
                {
                    "id": 1,
                    "last_scan": datetime(2026, 7, 13, 12, 15),
                    "barcode": "LQC-001",
                    "serial": "S-001",
                },
                {
                    "id": 2,
                    "last_scan": datetime(2026, 7, 13, 18, 15),
                    "barcode": "LQC-002",
                    "serial": "S-002",
                },
                {
                    "id": 3,
                    "last_scan": datetime(2026, 7, 13, 23, 15),
                    "barcode": "LQC-003",
                    "serial": "S-003",
                }
            ]

        def close(self):
            pass

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            pass

    connection = FakeConnection()
    monkeypatch.setattr(ai_reports, "report_allowed", lambda *_: True)
    monkeypatch.setattr(
        ai_reports, "now_local", lambda: datetime(2026, 7, 13, 20, 46, 31)
    )
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)

    result = ai_reports.run_report("ana", "quality_lqc", {})

    select_sql, select_params = connection.cursor_instance.executions[-1]
    assert "`last_scan` >= %s" in select_sql
    assert "`last_scan` <= %s" in select_sql
    assert select_params[:2] == (
        "2026-07-13 07:30:00",
        "2026-07-14 07:29:59",
    )
    assert result["filters"]["date_scope"] == "today"
    assert result["filters"]["operational_date"] == "2026-07-13"
    assert result["summary"]["operational_window"] == "07:30 a 07:30 del día siguiente"
    assert [item["count"] for item in result["summary"]["by_shift"]] == [1, 1, 1]
    assert [row["turno"] for row in result["rows"]] == [
        "DIA",
        "TIEMPO EXTRA",
        "NOCHE",
    ]


def test_lqc_respeta_y_normaliza_rango_explicito():
    filters = ai_reports._quality_lqc_date_filters(
        {"date_from": "2026-07-01", "date_to": "2026-07-03"}
    )
    assert filters == {
        "date_from": "2026-07-01 07:30:00",
        "date_to": "2026-07-04 07:29:59",
    }
    one_day = ai_reports._quality_lqc_date_filters({"date_from": "2026-07-13"})
    assert one_day == {
        "date_from": "2026-07-13 07:30:00",
        "date_to": "2026-07-14 07:29:59",
    }


def test_tool_schema_es_cerrado(monkeypatch):
    monkeypatch.setattr(
        ai_reports,
        "allowed_reports",
        lambda _username: [{"key": "bom", "title": "BOM", "description": "x"}],
    )
    schema = ai_reports.query_tool_schema("ana")
    assert schema["strict"] is True
    assert schema["parameters"]["additionalProperties"] is False
    assert schema["parameters"]["properties"]["report"]["enum"] == ["bom"]


def test_raw_es_fuente_maestra_autorizada_para_ct_y_uph():
    report = ai_reports.REPORTS["raw_model_standards"]
    assert report.source == "raw"
    assert report.permissions == (
        ("LISTA_INFORMACIONBASICA", "Control de produccion", "Control de modelos"),
    )
    assert "part_no" in report.description
    assert "uph" in report.description.lower()


def test_instrucciones_priorizan_raw_para_uph_y_ct():
    instructions = ai_openai.build_instructions({"language": "es"})
    assert "consulta primero raw_model_standards" in instructions
    assert "No calcules ni infieras UPH" in instructions
    assert "7421 debe encontrar EBR80757421" in instructions


def test_instrucciones_exigen_proceso_actual_al_replanear_hoy():
    instructions = ai_openai.build_instructions(
        {"language": "es", "plan_tools_enabled": True}
    )
    assert "¿en qué proceso o lote van las líneas?" in instructions
    assert "proceso_actual" in instructions


def test_instrucciones_lqc_explican_fecha_de_hoy():
    instructions = ai_openai.build_instructions(
        {
            "language": "es",
            "current_local_datetime": "2026-07-13 20:46:31",
        }
    )
    assert "resultados de hoy" in instructions
    assert "operational_date" in instructions
    assert "no digas \"sin filtros\"" in instructions
    assert "Día 07:30–17:30" in instructions
    assert "Tiempo extra 17:30–22:00" in instructions
    assert "Noche 22:00–07:30" in instructions
    assert "conteo por turno y actividad por hora" in instructions
    assert "2026-07-13 20:46:31" in instructions


def test_reporte_estado_lineas_combina_areas_autorizadas(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.sql = ""
            self.executions = []

        def execute(self, sql, params):
            self.sql = sql
            self.executions.append((sql, params))

        def fetchall(self):
            if "`plan_main`" in self.sql:
                return [
                    {
                        "lot_no": "ASSY-1",
                        "working_date": "2026-07-13",
                        "line": "M4",
                        "shift": None,
                        "model_code": "EBR80757421",
                        "part_no": "ABQ1",
                        "project": "QT-T",
                        "process": "ASSY",
                        "plan_count": 100,
                        "produced_count": 60,
                        "output": 55,
                        "delivered_count": 50,
                        "status": "EN PROGRESO",
                    }
                ]
            if "`plan_smt`" in self.sql:
                return [
                    {
                        "lot_no": "SMT-1",
                        "working_date": "2026-07-13",
                        "line": "S1",
                        "shift": "DIA",
                        "model_code": "EBR80757421",
                        "part_no": "EBR80757421",
                        "project": "QT-T",
                        "process": "SMT",
                        "plan_count": 200,
                        "produced_count": 100,
                        "output": 90,
                        "delivered_count": 80,
                        "status": "EN PROGRESO",
                    }
                ]
            return []

        def close(self):
            pass

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            pass

    connection = FakeConnection()
    monkeypatch.setattr(
        ai_reports,
        "puede_boton",
        lambda _username, _page, _section, button: button
        in {"Control de produccion ASSY", "Control de produccion SMT Plan"},
    )
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)

    result = ai_reports.run_report(
        "ana",
        "line_status_today",
        {"date_from": "2026-07-13", "date_to": "2026-07-13"},
    )

    assert result["sources"] == ["plan_main", "plan_smt"]
    assert result["queried_areas"] == ["ASSY", "SMT"]
    assert result["omitted_areas"] == ["IMT"]
    assert result["summary"]["overall"]["plan_count"] == 300
    assert result["summary"]["overall"]["produced_count"] == 160
    assert result["summary"]["overall"]["produced_progress_pct"] == 53.33
    assert result["summary"]["by_area"]["ASSY"]["output_progress_pct"] == 55
    assert result["visualization"]["type"] == "production_area_progress"
    assert result["visualization"]["areas"][0] == {
        "area": "ASSY",
        "goal": 100.0,
        "produced": 60.0,
        "output": 55.0,
        "produced_pct": 60.0,
        "output_pct": 55.0,
        "active_lines": 1,
        "plans": 1,
        "status_counts": {"EN PROGRESO": 1},
    }
    assert all(execution[1][:2] == ("2026-07-13", "2026-07-13") for execution in connection.cursor_instance.executions)


def test_instrucciones_priorizan_control_produccion_para_estado_lineas():
    instructions = ai_openai.build_instructions({"language": "es"})
    assert "consulta primero line_status_today" in instructions
    assert "Así van las líneas hoy" in instructions
    assert "respuesta ejecutiva y atractiva" in instructions
    assert "no repitas el nombre de un archivo adjunto" in instructions
    assert "No uses production_plans ni lg_plan_daily" in instructions
    assert "plan_main" in instructions
    assert "plan_imd" in instructions
    assert "plan_smt" in instructions


def test_reporte_estado_calidad_usa_lqc_ict_liberacion_y_vision(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.sql = ""
            self.executions = []

        def execute(self, sql, params=None):
            self.sql = " ".join(str(sql).split())
            self.executions.append((self.sql, params))

        def fetchone(self):
            if "AS inspected FROM box_scans" in self.sql:
                return {"inspected": 1000}
            if "AS defects FROM defect_data" in self.sql:
                return {"defects": 10}
            if "control_calidad_ppm_targets" in self.sql:
                return {"target_ppm": 15000}
            if "FROM history_ict" in self.sql:
                return {"total": 100, "unique_units": 90, "passed": 92, "failed": 8}
            if "FROM box_scans" in self.sql:
                return {"total": 80, "unique_units": 75, "lots": 4}
            if "FROM history_vision" in self.sql:
                return {"total": 50, "unique_units": 48, "passed": 47, "failed": 3}
            return {}

        def close(self):
            pass

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            pass

    connection = FakeConnection()
    allowed_buttons = {
        "PPM's LQC",
        "Historial de maquina ICT",
        "Historial de liberacion LQC",
        "Historial de maquina vision",
    }
    monkeypatch.setattr(
        ai_reports,
        "puede_boton",
        lambda _username, _page, _section, button: button in allowed_buttons,
    )
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)

    result = ai_reports.run_report(
        "ana",
        "quality_status_today",
        {"date_from": "2026-07-13", "date_to": "2026-07-13"},
    )

    assert result["queried_sources"] == [
        "lqc_results", "ict_history", "lqc_release_history", "vision_history"
    ]
    assert result["omitted_sources"] == []
    assert result["summary"]["lqc_results"]["ppm"] == 10000
    assert result["summary"]["lqc_results"]["within_target"] is True
    assert result["summary"]["ict_history"]["pass_rate_pct"] == 92
    assert result["summary"]["lqc_release_history"]["duplicates"] == 5
    assert result["summary"]["vision_history"]["pass_rate_pct"] == 94
    assert result["visualization"]["type"] == "quality_today_overview"
    assert "AOI" not in result["source"]
    assert result["truncated"] is False
    assert all(params is not None for _, params in connection.cursor_instance.executions)


def test_instrucciones_definen_alcance_correcto_de_calidad():
    instructions = ai_openai.build_instructions({"language": "es"})
    assert "consulta primero y únicamente quality_status_today" in instructions
    assert "Resultados LQC" in instructions
    assert "Historial ICT" in instructions
    assert "Historial de liberación LQC" in instructions
    assert "Historial Vision" in instructions
    assert "AOI no forma parte" in instructions
    assert "No generes ni anuncies Excel" in instructions


def test_safety_identifier_no_expone_username(monkeypatch):
    monkeypatch.setenv("SECRET_KEY", "test-secret")
    identifier = ai_openai.safety_identifier("nombre.apellido@ilsan.com")
    assert identifier.startswith("mes_")
    assert "nombre" not in identifier
    assert identifier == ai_openai.safety_identifier("nombre.apellido@ilsan.com")


def test_stream_openai_usa_store_false_y_gpt55(monkeypatch):
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setenv("SECRET_KEY", "test-secret")
    monkeypatch.setenv("OPENAI_MODEL", "gpt-5.5")
    captured = {}
    response = SimpleNamespace(
        id="resp_1",
        output=[],
        usage=SimpleNamespace(input_tokens=10, output_tokens=3, total_tokens=13),
    )

    class Responses:
        def create(self, **kwargs):
            captured.update(kwargs)
            return iter(
                [
                    SimpleNamespace(type="response.output_text.delta", delta="Hola"),
                    SimpleNamespace(type="response.completed", response=response),
                ]
            )

    monkeypatch.setattr(ai_openai, "_client", lambda: SimpleNamespace(responses=Responses()))
    events = list(
        ai_openai.stream_response(
            username="ana",
            context={"language": "es"},
            messages=[{"role": "user", "content": "Hola"}],
            tools=[],
            execute_tool=lambda *_: {},
        )
    )
    assert captured["model"] == "gpt-5.5"
    assert captured["store"] is False
    assert captured["stream"] is True
    assert captured["safety_identifier"].startswith("mes_")
    assert events[0] == {"event": "delta", "data": {"text": "Hola"}}
    assert events[-1]["event"] == "usage"


def test_stream_openai_emite_visualizaciones_de_herramientas(monkeypatch):
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")
    monkeypatch.setenv("SECRET_KEY", "test-secret")
    responses = [
        SimpleNamespace(
            id="resp_tool",
            output=[
                SimpleNamespace(
                    type="function_call",
                    id="fc_1",
                    call_id="call_1",
                    name="query_mes_report",
                    arguments="{}",
                    status="completed",
                )
            ],
            usage=SimpleNamespace(input_tokens=5, output_tokens=1, total_tokens=6),
        ),
        SimpleNamespace(
            id="resp_final",
            output=[],
            usage=SimpleNamespace(input_tokens=4, output_tokens=2, total_tokens=6),
        ),
    ]

    class Responses:
        def create(self, **_kwargs):
            response = responses.pop(0)
            return iter([SimpleNamespace(type="response.completed", response=response)])

    visualization = {"id": "viz-1", "type": "production_area_progress", "areas": []}
    artifact = {"id": "file-1", "type": "xlsx"}
    monkeypatch.setattr(ai_openai, "_client", lambda: SimpleNamespace(responses=Responses()))
    events = list(
        ai_openai.stream_response(
            username="ana",
            context={"language": "es"},
            messages=[{"role": "user", "content": "¿Cómo van las líneas?"}],
            tools=[{"type": "function"}],
            execute_tool=lambda *_args: {
                "model_output": {"success": True},
                "public_summary": {},
                "client_events": [{"event": "visualization", "data": visualization}],
                "client_event": {"event": "artifact_ready", "data": artifact},
            },
        )
    )
    assert {"event": "visualization", "data": visualization} in events
    assert events.index({"event": "visualization", "data": visualization}) < events.index(
        {"event": "artifact_ready", "data": artifact}
    )


def test_estado_lineas_no_exporta_excel_solo_por_cantidad_de_filas():
    result = {"row_count": 80, "truncated": False}
    assert not ai_assistant._should_auto_export_report(
        "line_status_today", result, large_request=False, threshold=50
    )
    assert not ai_assistant._should_auto_export_report(
        "line_status_today", result, large_request=True, threshold=50
    )
    assert ai_assistant._should_auto_export_report(
        "line_status_today", {**result, "truncated": True}, large_request=False, threshold=50
    )
    assert ai_assistant._should_auto_export_report(
        "raw_model_standards", result, large_request=False, threshold=50
    )
    assert not ai_assistant._should_auto_export_report(
        "quality_status_today", result, large_request=True, threshold=50
    )


def test_panel_renderiza_y_recupera_grafica_de_avance():
    root = Path(__file__).resolve().parents[1]
    javascript = (root / "app/static/js/ai-assistant.js").read_text(encoding="utf-8")
    stylesheet = (root / "app/static/css/ai-assistant.css").read_text(encoding="utf-8")
    assert "appendVisualization(visualization" in javascript
    assert "message.content_json?.visualizations" in javascript
    assert "event === 'visualization'" in javascript
    assert "renderMarkdown(bubble, source)" in javascript
    assert "appendInlineMarkdown(parent, source)" in javascript
    assert "this.renderMarkdown(this.activeAssistantBubble, this.activeAssistantText)" in javascript
    assert "startThinkingIndicator()" in javascript
    assert "showThinkingIndicator(this.t('reasoning'))" in javascript
    assert "reasoning:'Razonando…'" in javascript
    assert ".ai-production-bar-fill.produced" in stylesheet
    assert ".ai-production-bar-fill.output" in stylesheet
    assert ".ai-bubble.ai-markdown" in stylesheet
    assert ".ai-markdown-table-wrap" in stylesheet
    assert ".ai-thinking-indicator" in stylesheet
    assert "@keyframes ai-thinking-dot" in stylesheet
    assert "appendQualityVisualization(visualization" in javascript
    assert "quality_today_overview" in javascript
    assert ".ai-quality-overview" in stylesheet
    assert "setExpanded(expanded, persist = true)" in javascript
    assert "mes-ai-expanded" in javascript
    assert "this.panel.classList.toggle('expanded'" in javascript
    assert ".ai-panel.expanded" in stylesheet
    assert ".ai-panel.expanded {" in stylesheet
    assert "width:100vw;" in stylesheet
    assert "grid-template-columns:clamp(245px,18vw,290px) minmax(0,1fr)" in stylesheet
    assert "--ai-chat-column:860px" in stylesheet
    assert "calc((100% - var(--ai-chat-column))/2)" in stylesheet
    assert ".ai-expanded-sidebar" in stylesheet
    assert ".ai-sidebar-conversation-list" in stylesheet
    assert "querySelectorAll('[data-ai-conversation-list]')" in javascript
    assert "highlightConversation()" in javascript
    assert "bindLauncherDrag()" in javascript
    assert "mes-ai-launcher-side" in javascript
    assert "setPointerCapture(event.pointerId)" in javascript
    assert "event.clientX < window.innerWidth / 2" in javascript
    assert '.ai-launcher[data-side="left"]' in stylesheet
    assert ".ai-launcher.dragging" in stylesheet
    assert ".ai-panel-brand" in stylesheet
    assert ".ai-expand-button" in stylesheet


def test_ddl_incluye_tablas_y_permisos(monkeypatch):
    executed = []

    class Cursor:
        def execute(self, sql, params=None):
            executed.append((sql, params))

        def fetchone(self):
            return None

        def close(self):
            pass

    class Connection:
        def cursor(self):
            return Cursor()

        def commit(self):
            pass

        def rollback(self):
            pass

        def close(self):
            pass

    monkeypatch.setattr(ai_store, "get_db_connection", lambda: Connection())
    ai_store.init_ai_assistant_tables()
    ddl = "\n".join(sql for sql, _ in executed)
    for table in (
        "ai_conversations", "ai_messages", "ai_tool_executions",
        "ai_usage_daily", "ai_usage_limits", "ai_knowledge_documents", "ai_artifacts",
    ):
        assert table in ddl
    permission_params = [params for _, params in executed if params and ai_store.AI_PAGE in params]
    assert {params[2] for params in permission_params} == {
        ai_store.AI_PERMISSION_USE,
        ai_store.AI_PERMISSION_ARTIFACTS,
        ai_store.AI_PERMISSION_AUDIT,
        ai_store.AI_PERMISSION_LIMITS,
    }


def test_main_template_incluye_panel_y_metadatos_de_permiso():
    root = Path(__file__).resolve().parents[1]
    main = (root / "app/templates/MainTemplate.html").read_text(encoding="utf-8")
    partial = (root / "app/templates/components/ai_assistant.html").read_text(encoding="utf-8")
    assert "components/ai_assistant.html" in main
    assert "ai-assistant.js" in main
    assert 'data-permiso-pagina="MAIN_TEMPLATE"' in partial
    assert 'data-permiso-boton="Usar asistente IA"' in partial
    assert 'id="ai-expand"' in partial
    assert 'aria-pressed="false"' in partial
    assert 'class="ai-workspace"' in partial
    assert 'data-ai-conversation-list' in partial
    assert 'data-ai-action="new-chat"' in partial
    assert 'id="ai-delete-chat"' in partial
    assert 'class="ai-launcher-logo"' in partial
    assert "icons/1538298822.svg" in partial
    assert '>Asistente</span>' not in partial.split('</button>', 1)[0]
    assert "20260715a" in main


def test_cliente_permite_eliminar_chat_con_confirmacion():
    root = Path(__file__).resolve().parents[1]
    javascript = (root / "app/static/js/ai-assistant.js").read_text(encoding="utf-8")
    stylesheet = (root / "app/static/css/ai-assistant.css").read_text(encoding="utf-8")

    assert "deleteConversation(conversation = this.currentConversation)" in javascript
    assert "window.confirm" in javascript
    assert "{method:'DELETE'}" in javascript
    assert "ai-list-delete" in javascript
    assert ".ai-list-delete" in stylesheet


def test_eliminar_conversacion_borra_registros_y_archivos_privados(tmp_path, monkeypatch):
    artifact_dir = tmp_path / "ai_artifacts"
    artifact_dir.mkdir()
    inside = artifact_dir / "inside.xlsx"
    outside = tmp_path / "outside.xlsx"
    inside.write_bytes(b"inside")
    outside.write_bytes(b"outside")

    class FakeCursor:
        def __init__(self):
            self.mode = ""
            self.rowcount = 0
            self.executions = []

        def execute(self, sql, params=None):
            self.executions.append((sql, params))
            normalized = " ".join(sql.split())
            if normalized.startswith("SELECT id, public_id, title"):
                self.mode = "conversation"
                self.rowcount = 0
            elif normalized.startswith("SELECT storage_path"):
                self.mode = "artifacts"
                self.rowcount = 0
            elif "DELETE FROM ai_tool_executions" in normalized:
                self.rowcount = 4
            elif "DELETE FROM ai_artifacts" in normalized:
                self.rowcount = 2
            elif "DELETE FROM ai_messages" in normalized:
                self.rowcount = 7
            elif "DELETE FROM ai_conversations" in normalized:
                self.rowcount = 1

        def fetchone(self):
            if self.mode == "conversation":
                return {"id": 19, "public_id": "chat-19", "title": "Prueba"}
            return None

        def fetchall(self):
            if self.mode == "artifacts":
                return [
                    {"storage_path": str(inside)},
                    {"storage_path": str(outside)},
                ]
            return []

        def close(self):
            pass

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()
            self.committed = False
            self.rolled_back = False

        def cursor(self):
            return self.cursor_instance

        def commit(self):
            self.committed = True

        def rollback(self):
            self.rolled_back = True

        def close(self):
            pass

    connection = FakeConnection()
    monkeypatch.setattr(ai_store, "get_db_connection", lambda: connection)
    monkeypatch.setattr(ai_store, "artifact_root", lambda: artifact_dir.resolve())

    deleted = ai_store.delete_conversation("chat-19", "ana")

    assert connection.committed is True
    assert deleted["messages"] == 7
    assert deleted["tool_executions"] == 4
    assert deleted["artifacts"] == 2
    assert deleted["files_removed"] == 1
    assert deleted["files_failed"] == 1
    assert not inside.exists()
    assert outside.exists()


def test_bootstrap_requiere_sesion(client):
    response = client.get(
        "/api/ai/bootstrap",
        headers={"Content-Type": "application/json", "Accept": "application/json"},
    )
    assert response.status_code == 401


def test_bootstrap_respeta_permiso_y_superadmin(client, monkeypatch):
    from app.api.portal import ai_assistant
    from app.api.shared import permisos

    class FakeAuth:
        role = "operador"

        def obtener_rol_principal_usuario(self, _username):
            return self.role

        def verificar_permiso_boton(self, *_args):
            return False

    fake = FakeAuth()
    monkeypatch.setattr(permisos, "_auth", lambda: fake)
    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["operador"]

    response = client.get("/api/ai/bootstrap", headers={"Content-Type": "application/json"})
    assert response.status_code == 403

    fake.role = "superadmin"
    monkeypatch.setattr(ai_assistant, "allowed_reports", lambda _username: [])
    monkeypatch.setattr(ai_assistant, "get_usage", lambda *_: {})
    monkeypatch.setattr(ai_assistant, "effective_limits", lambda *_: {})
    monkeypatch.setattr(ai_assistant, "permisos_botones", lambda *_: {})
    response = client.get("/api/ai/bootstrap", headers={"Content-Type": "application/json"})
    assert response.status_code == 200


def test_conversacion_ajena_responde_404(client, monkeypatch):
    from app.api.portal import ai_assistant
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(
        ai_assistant,
        "get_conversation",
        lambda _public_id: {"id": 9, "public_id": "other", "username": "otro"},
    )
    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]
    response = client.get(
        "/api/ai/conversations/other/messages",
        headers={"Content-Type": "application/json"},
    )
    assert response.status_code == 404
    delete_response = client.delete(
        "/api/ai/conversations/other",
        headers={"Content-Type": "application/json"},
    )
    assert delete_response.status_code == 404


def test_propietario_puede_eliminar_conversacion(client, monkeypatch):
    from app.api.shared import permisos

    class SuperadminAuth:
        def obtener_rol_principal_usuario(self, _username):
            return "superadmin"

    captured = {}
    monkeypatch.setattr(permisos, "_auth", lambda: SuperadminAuth())
    monkeypatch.setattr(
        ai_assistant,
        "get_conversation",
        lambda _public_id: {"id": 12, "public_id": "mine", "username": "ana"},
    )
    monkeypatch.setattr(
        ai_assistant,
        "delete_conversation",
        lambda public_id, username: {
            "public_id": public_id,
            "messages": 5,
            "tool_executions": 2,
            "artifacts": 1,
            "files_removed": 1,
            "files_failed": 0,
        },
    )
    monkeypatch.setattr(
        ai_assistant,
        "_audit",
        lambda action, description, details=None, result="EXITOSO": captured.update(
            {"action": action, "description": description, "details": details}
        ),
    )
    with client.session_transaction() as sess:
        sess["usuario"] = "ana"
        sess["roles"] = ["superadmin"]

    response = client.delete(
        "/api/ai/conversations/mine",
        headers={"Content-Type": "application/json"},
    )

    assert response.status_code == 200
    assert response.get_json()["deleted"]["messages"] == 5
    assert captured["action"] == "ELIMINAR_CHAT"
    assert captured["details"]["files_removed"] == 1


def test_almacen_turno_nocturno_usa_solo_la_noche_mas_reciente(monkeypatch):
    monkeypatch.setattr(
        ai_reports,
        "now_local",
        lambda: datetime(2026, 7, 14, 7, 41, 1),
    )

    window = ai_reports._warehouse_shift_window({"shift": "noche"})

    assert window["scope"] == "latest_shift"
    assert window["window_start"] == datetime(2026, 7, 13, 22, 0, 0)
    assert window["window_end"] == datetime(2026, 7, 14, 7, 30, 0)


def test_almacen_sin_fechas_usa_turno_actual_hasta_ahora(monkeypatch):
    monkeypatch.setattr(
        ai_reports,
        "now_local",
        lambda: datetime(2026, 7, 14, 7, 41, 1),
    )

    window = ai_reports._warehouse_shift_window({})

    assert window["scope"] == "current_shift"
    assert window["shift"] == "dia"
    assert window["window_start"] == datetime(2026, 7, 14, 7, 30, 0)
    assert window["window_end"] == datetime(2026, 7, 14, 7, 41, 1)


def test_conteo_almacen_consulta_historial_entradas_en_una_sola_agregacion(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.executions = []

        def execute(self, sql, params=None):
            self.executions.append((sql, params))

        def fetchone(self):
            return {
                "registros": 37,
                "movimientos_unicos": 35,
                "numeros_parte": 18,
                "cantidad_total": 1240,
            }

        def close(self):
            return None

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            return None

    connection = FakeConnection()
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)
    monkeypatch.setattr(ai_reports, "puede_boton", lambda *_args: True)
    monkeypatch.setattr(
        ai_reports,
        "now_local",
        lambda: datetime(2026, 7, 14, 7, 41, 1),
    )

    result = ai_reports.run_report(
        "ana",
        "warehouse_shift_activity",
        {"movement_type": "entradas", "shift": "noche"},
    )

    assert len(connection.cursor_instance.executions) == 1
    sql, params = connection.cursor_instance.executions[0]
    assert "FROM `control_material_almacen`" in sql
    assert "`fecha_recibo` >= %s" in sql
    assert "embarques" not in sql.lower()
    assert params == ("2026-07-13 22:00:00", "2026-07-14 07:30:00")
    assert result["summary"]["counts"] == {"entradas": 37}
    assert result["source"] == "Control de material"


def test_contexto_almacen_prioriza_control_material_y_respuesta_corta():
    instructions = ai_openai.build_instructions(
        {"language": "es", "current_local_datetime": "2026-07-14 07:41:01"}
    )

    assert '"almacén" significa por defecto Información básica > Control de material' in instructions
    assert "NO significa Almacén de Embarques" in instructions
    assert "consulta primero y únicamente warehouse_shift_activity" in instructions
    assert "Nunca interpretes \"turno nocturno\" como todos" in instructions
    assert "máximo tres oraciones" in instructions
    for module in (
        "Historial de entradas",
        "Historial de salidas",
        "Historial de retornos",
        "Inventario actual",
        "Facturas / Invoice",
        "Valorización de inventario",
        "Lista de compras",
    ):
        assert module in instructions


def test_schema_incluye_turno_y_tipo_movimiento_canonicos(monkeypatch):
    monkeypatch.setattr(
        ai_reports,
        "allowed_reports",
        lambda _username: [{"key": "warehouse_shift_activity"}],
    )

    schema = ai_reports.query_tool_schema("ana")
    filters = schema["parameters"]["properties"]["filters"]

    assert filters["properties"]["shift"]["enum"] == [
        "actual", "dia", "tiempo_extra", "noche", None,
    ]
    assert filters["properties"]["movement_type"]["enum"] == [
        "entradas", "salidas", "retornos", "entradas_salidas", "todos", None,
    ]
    assert "shift" in filters["required"]
    assert "movement_type" in filters["required"]


def test_consulta_simple_almacen_activa_modo_compacto_sin_archivos():
    assert ai_assistant._compact_warehouse_count_request(
        "En almacén cuántas entradas tuvo en el turno nocturno?"
    )
    assert ai_assistant._compact_warehouse_count_request(
        "How many warehouse entries were recorded on the night shift?"
    )
    assert not ai_assistant._compact_warehouse_count_request(
        "Dame en Excel las entradas de almacén del turno nocturno"
    )
    assert not ai_assistant._compact_warehouse_count_request(
        "Cuántos embarques hubo en el turno nocturno"
    )


def test_schema_puede_restringirse_al_reporte_compacto_de_almacen(monkeypatch):
    monkeypatch.setattr(
        ai_reports,
        "allowed_reports",
        lambda _username: [
            {"key": "warehouse_shift_activity"},
            {"key": "shipping_inventory"},
            {"key": "material_entries"},
        ],
    )

    schema = ai_reports.query_tool_schema("ana", {"warehouse_shift_activity"})

    assert schema["parameters"]["properties"]["report"]["enum"] == [
        "warehouse_shift_activity"
    ]


def test_analisis_almacen_agrupa_entradas_y_salidas_por_parte_y_turno(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.executions = []
            self.sql = ""

        def execute(self, sql, params=None):
            self.sql = sql
            self.executions.append((sql, params))

        def fetchall(self):
            if "control_material_almacen" in self.sql:
                return [
                    {
                        "tipo_movimiento": "entradas",
                        "numero_parte": "MAT-100",
                        "turno": "DIA",
                        "registros": 3,
                        "movimientos_unicos": 3,
                        "cantidad_total": 120,
                        "cancelados": 0,
                        "primera_actividad": datetime(2026, 7, 2, 8, 0),
                        "ultima_actividad": datetime(2026, 7, 8, 9, 0),
                    }
                ]
            return [
                {
                    "tipo_movimiento": "salidas",
                    "numero_parte": "MAT-100",
                    "turno": "NOCHE",
                    "registros": 2,
                    "movimientos_unicos": 2,
                    "cantidad_total": 45,
                    "cancelados": 0,
                    "primera_actividad": datetime(2026, 7, 3, 23, 0),
                    "ultima_actividad": datetime(2026, 7, 9, 2, 0),
                }
            ]

        def close(self):
            return None

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            return None

    connection = FakeConnection()
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)
    monkeypatch.setattr(ai_reports, "puede_boton", lambda *_args: True)
    monkeypatch.setattr(
        ai_reports,
        "now_local",
        lambda: datetime(2026, 7, 14, 8, 3, 13),
    )

    result = ai_reports.run_report(
        "ana",
        "warehouse_analysis",
        {"movement_type": "entradas_salidas"},
        for_artifact=True,
        limit=10000,
    )

    assert len(connection.cursor_instance.executions) == 2
    assert all("GROUP BY numero_parte, turno" in sql for sql, _ in connection.cursor_instance.executions)
    assert all("shipping" not in sql.lower() for sql, _ in connection.cursor_instance.executions)
    first_params = connection.cursor_instance.executions[0][1]
    assert first_params[1:3] == ("2026-07-01 00:00:00", "2026-07-14 08:03:13")
    assert result["filters"]["date_scope"] == "current_month"
    assert result["summary"]["total_registros"] == 5
    assert result["summary"]["cantidad_total"] == 165
    assert result["summary"]["numeros_parte"] == 1
    assert {row["turno"] for row in result["rows"]} == {"DIA", "NOCHE"}


def test_analisis_lqc_reutiliza_parte_linea_y_turno_de_api_canonica(monkeypatch):
    class FakeCursor:
        def __init__(self):
            self.execution = None

        def execute(self, sql, params=None):
            self.execution = (sql, params)

        def fetchall(self):
            return [
                {
                    "numero_parte": "EBR80757421",
                    "linea": "M4",
                    "turno": "DIA",
                    "cantidad_total": 250,
                    "unidades_unicas": 248,
                    "escaneos_repetidos": 2,
                    "lotes": 3,
                    "primera_actividad": datetime(2026, 7, 13, 8, 0),
                    "ultima_actividad": datetime(2026, 7, 13, 16, 30),
                }
            ]

        def close(self):
            return None

    class FakeConnection:
        def __init__(self):
            self.cursor_instance = FakeCursor()

        def cursor(self):
            return self.cursor_instance

        def close(self):
            return None

    connection = FakeConnection()
    monkeypatch.setattr(ai_reports, "get_db_connection", lambda: connection)
    monkeypatch.setattr(ai_reports, "puede_boton", lambda *_args: True)

    result = ai_reports.run_report(
        "ana",
        "quality_lqc_analysis",
        {"date_from": "2026-07-13", "date_to": "2026-07-13"},
        for_artifact=True,
    )

    sql, params = connection.cursor_instance.execution
    assert "FROM box_scans b" in sql
    assert "LEFT JOIN plan_main p ON p.lot_no = b.lot_no" in sql
    assert "p.part_no" in sql
    assert "GROUP BY numero_parte, linea, turno" in sql
    assert params[:2] == ("2026-07-13 07:30:00", "2026-07-14 07:29:59")
    assert result["summary"]["detail_level"] == "numero_parte_linea_turno"
    assert result["summary"]["cantidad_total"] == 250


def test_excel_analitico_almacen_separa_movimientos_y_conserva_detalle(tmp_path):
    result = {
        "report": "warehouse_analysis",
        "title": "Análisis detallado de almacén",
        "source": "Control de material",
        "filters": {
            "date_from": "2026-07-01 00:00:00",
            "date_to": "2026-07-14 08:03:13",
            "movement_type": "entradas_salidas",
        },
        "columns": [
            "tipo_movimiento", "numero_parte", "turno", "registros",
            "movimientos_unicos", "cantidad_total", "cancelados",
        ],
        "rows": [
            {
                "tipo_movimiento": "entradas", "numero_parte": "MAT-100",
                "turno": "DIA", "registros": 3, "movimientos_unicos": 3,
                "cantidad_total": 120, "cancelados": 0,
            },
            {
                "tipo_movimiento": "salidas", "numero_parte": "MAT-100",
                "turno": "NOCHE", "registros": 2, "movimientos_unicos": 2,
                "cantidad_total": 45, "cancelados": 0,
            },
        ],
        "row_count": 2,
        "truncated": False,
        "summary": {
            "total_registros": 5,
            "cantidad_total": 165,
            "numeros_parte": 1,
            "detail_level": "numero_parte_turno",
            "by_shift": [
                {"label": "DIA", "registros": 3, "cantidad_total": 120},
                {"label": "NOCHE", "registros": 2, "cantidad_total": 45},
            ],
            "by_movement": [
                {"label": "entradas", "registros": 3, "cantidad_total": 120},
                {"label": "salidas", "registros": 2, "cantidad_total": 45},
            ],
            "top_parts": [
                {"label": "MAT-100", "registros": 5, "cantidad_total": 165},
            ],
        },
    }
    target = tmp_path / "analisis_almacen.xlsx"

    ai_artifacts._build_excel(
        result,
        "Análisis de almacén",
        "es",
        target,
        include_summary=True,
        include_charts=True,
    )

    wb = load_workbook(target, data_only=False)
    assert wb.sheetnames == [
        "Resumen", "Analisis por parte", "Entradas", "Salidas", "Criterios"
    ]
    assert wb["Analisis por parte"]["B2"].value == "MAT-100"
    assert wb["Entradas"].max_row == 2
    assert wb["Salidas"].max_row == 2
    assert len(wb["Resumen"]._charts) == 2


def test_instrucciones_distinguen_conteo_de_analisis_detallado():
    instructions = ai_openai.build_instructions(
        {"language": "es", "current_local_datetime": "2026-07-14 08:03:13"}
    )

    assert "La palabra \"análisis\"" in instructions
    assert "usa warehouse_analysis" in instructions
    assert "movement_type=entradas_salidas" in instructions
    assert "No uses warehouse_shift_activity para un análisis" in instructions
    assert "usa quality_lqc_analysis" in instructions
    assert "número de parte, línea y turno" in instructions


def test_enrutador_de_analisis_selecciona_reporte_detallado():
    assert ai_assistant._detailed_analysis_report(
        "Haz un análisis de todas las entradas y salidas de almacén este mes en Excel"
    ) == "warehouse_analysis"
    assert ai_assistant._detailed_analysis_report(
        "Analiza LQC por número de parte y turno"
    ) == "quality_lqc_analysis"
    assert ai_assistant._detailed_analysis_report(
        "Cuántas entradas tuvo almacén en el turno nocturno"
    ) is None
